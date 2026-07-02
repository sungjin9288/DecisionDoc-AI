"""app/services/attachment/format_extractors.py — non-PDF format extractors.

Covers office documents (DOCX, PPTX, HWPX, Excel), structured text formats
(CSV/TSV, JSON, YAML, XML, HTML, RTF), OpenDocument formats (ODT/ODS/ODP),
and plain text. ZIP archive traversal lives in ``core`` alongside
``_extract_by_extension`` since the two are mutually recursive.
"""
from __future__ import annotations

import csv
import html
import io
import json
import re
import zipfile

from defusedxml import ElementTree as SafeET

from app.services.attachment.constants import MAX_CHARS_PER_FILE, AttachmentError


def _extract_docx(raw: bytes, filename: str) -> str:
    """Extract text from DOCX including tables, preserving document order."""
    try:
        from docx import Document
        from docx.table import Table
        from docx.text.paragraph import Paragraph
    except ImportError as exc:
        raise AttachmentError(
            "DOCX 파싱 라이브러리가 설치되어 있지 않습니다. "
            "'pip install python-docx'를 실행하세요."
        ) from exc

    doc = Document(io.BytesIO(raw))
    sections: list[str] = []

    for block in doc.element.body:
        tag = block.tag.split("}")[-1] if "}" in block.tag else block.tag

        if tag == "p":
            para = Paragraph(block, doc)
            text = para.text.strip()
            if text:
                sections.append(text)

        elif tag == "tbl":
            table = Table(block, doc)
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                if any(cells):
                    rows.append(" | ".join(cells))
            if rows:
                sections.append("[표]\n" + "\n".join(rows))

    return "\n".join(sections)[:MAX_CHARS_PER_FILE]


def _extract_hwpx(raw: bytes, filename: str) -> str:
    """Extract text from HWPX (ZIP + XML format)."""
    try:
        with zipfile.ZipFile(io.BytesIO(raw)) as zf:
            section_files = sorted(
                f for f in zf.namelist()
                if re.match(r"Contents/section\d+\.xml", f)
            )
            if not section_files:
                raise AttachmentError(
                    f"{filename}: HWPX 섹션 파일을 찾을 수 없습니다"
                )
            texts: list[str] = []
            for sf in section_files:
                xml = zf.read(sf).decode("utf-8", errors="ignore")
                # <hh:t> or <t> tags carry the visible text
                matches = re.findall(
                    r"<(?:hh:)?t[^>]*>([^<]+)</(?:hh:)?t>", xml
                )
                texts.extend(m.strip() for m in matches if m.strip())

    except zipfile.BadZipFile as exc:
        raise AttachmentError(
            f"{filename}: 구형 바이너리 HWP는 직접 추출하지 못합니다. "
            "HWPX, PDF 또는 DOCX로 변환해 업로드하세요."
        ) from exc

    result = "\n".join(texts)
    if not result:
        raise AttachmentError(f"{filename}: 텍스트를 추출할 수 없습니다")
    return result[:MAX_CHARS_PER_FILE]


def _extract_excel(raw: bytes, filename: str) -> str:
    """Extract text from Excel (.xlsx / .xls) via openpyxl."""
    try:
        import openpyxl
    except ImportError as exc:
        raise AttachmentError(
            "Excel 파싱 라이브러리가 설치되어 있지 않습니다. "
            "'pip install openpyxl'를 실행하세요."
        ) from exc

    wb = openpyxl.load_workbook(
        io.BytesIO(raw), read_only=True, data_only=True
    )
    sheets: list[str] = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows: list[str] = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c).strip() if c is not None else "" for c in row]
            if any(c for c in cells if c):
                rows.append(" | ".join(cells))
        if rows:
            sheets.append(f"[시트: {sheet_name}]\n" + "\n".join(rows))
    wb.close()

    return "\n\n".join(sheets)[:MAX_CHARS_PER_FILE]


def _extract_csv(raw: bytes, *, delimiter: str) -> str:
    """Extract text from CSV."""
    text = raw.decode("utf-8", errors="ignore")
    reader = csv.reader(io.StringIO(text), delimiter=delimiter)
    rows = [" | ".join(row) for row in reader if any(row)]
    return "\n".join(rows)[:MAX_CHARS_PER_FILE]


def _extract_pptx(raw: bytes, filename: str) -> str:
    """Extract text from PPTX — slide title + body text in slide order."""
    try:
        from pptx import Presentation
        from pptx.util import Pt  # noqa: F401 — confirms pptx is available
    except ImportError as exc:
        raise AttachmentError(
            "PPTX 파싱 라이브러리가 설치되어 있지 않습니다. "
            "'pip install python-pptx'를 실행하세요."
        ) from exc

    prs = Presentation(io.BytesIO(raw))
    slides: list[str] = []

    for i, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = "".join(run.text for run in para.runs).strip()
                if text:
                    parts.append(text)
        if parts:
            slides.append(f"[슬라이드 {i}]\n" + "\n".join(parts))

    result = "\n\n".join(slides)
    if not result.strip():
        raise AttachmentError(f"{filename}: 텍스트를 추출할 수 없습니다")
    return result[:MAX_CHARS_PER_FILE]


def _extract_plain(raw: bytes) -> str:
    """Extract plain text (.txt / .md)."""
    try:
        return raw.decode("utf-8")[:MAX_CHARS_PER_FILE]
    except UnicodeDecodeError:
        return raw.decode("latin-1", errors="replace")[:MAX_CHARS_PER_FILE]


def _extract_json(raw: bytes, filename: str) -> str:
    text = _extract_plain(raw)
    try:
        parsed = json.loads(text)
    except json.JSONDecodeError:
        return text[:MAX_CHARS_PER_FILE]
    return json.dumps(parsed, ensure_ascii=False, indent=2)[:MAX_CHARS_PER_FILE]


def _extract_yaml(raw: bytes, filename: str) -> str:
    try:
        import yaml
    except ImportError as exc:
        raise AttachmentError(
            "YAML 파싱 라이브러리가 설치되어 있지 않습니다. 'pip install PyYAML'를 실행하세요."
        ) from exc

    text = _extract_plain(raw)
    try:
        parsed = yaml.safe_load(text)
    except yaml.YAMLError:
        return text[:MAX_CHARS_PER_FILE]
    if parsed is None:
        return text[:MAX_CHARS_PER_FILE]
    return yaml.safe_dump(parsed, allow_unicode=True, sort_keys=False)[:MAX_CHARS_PER_FILE]


def _extract_xml(raw: bytes, filename: str) -> str:
    text = _extract_plain(raw)
    try:
        root = SafeET.fromstring(text)
    except SafeET.ParseError:
        return text[:MAX_CHARS_PER_FILE]

    parts: list[str] = []
    for elem in root.iter():
        if elem.text and elem.text.strip():
            parts.append(elem.text.strip())
    result = "\n".join(parts)
    if not result:
        raise AttachmentError(f"{filename}: XML에서 텍스트를 추출할 수 없습니다")
    return result[:MAX_CHARS_PER_FILE]


def _extract_html(raw: bytes) -> str:
    text = _extract_plain(raw)
    text = re.sub(r"(?is)<(script|style)\b[^>]*>.*?</\1>", " ", text)
    text = re.sub(r"(?i)<br\s*/?>", "\n", text)
    text = re.sub(r"(?i)</(p|div|li|tr|h[1-6])>", "\n", text)
    text = re.sub(r"(?s)<[^>]+>", " ", text)
    text = html.unescape(text)
    text = re.sub(r"[ \t]+\n", "\n", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    text = re.sub(r"[ \t]{2,}", " ", text)
    return text.strip()[:MAX_CHARS_PER_FILE]


def _extract_rtf(raw: bytes) -> str:
    text = raw.decode("latin-1", errors="ignore")
    text = re.sub(
        r"\\'([0-9a-fA-F]{2})",
        lambda m: bytes.fromhex(m.group(1)).decode("cp1252", errors="ignore"),
        text,
    )
    text = re.sub(r"\\u(-?\d+)\??", lambda m: _decode_rtf_unicode(m.group(1)), text)
    text = text.replace("\\par", "\n").replace("\\line", "\n").replace("\\tab", "\t")
    text = re.sub(r"\\[a-zA-Z]+-?\d* ?", "", text)
    text = text.replace("{", "").replace("}", "")
    text = re.sub(r"\n{3,}", "\n\n", text)
    text = re.sub(r"[ \t]{2,}", " ", text)
    return text.strip()[:MAX_CHARS_PER_FILE]


def _decode_rtf_unicode(value: str) -> str:
    codepoint = int(value)
    if codepoint < 0:
        codepoint += 65536
    try:
        return chr(codepoint)
    except ValueError:
        return ""


def _extract_opendocument(raw: bytes, filename: str) -> str:
    try:
        with zipfile.ZipFile(io.BytesIO(raw)) as zf:
            if "content.xml" not in zf.namelist():
                raise AttachmentError(f"{filename}: OpenDocument content.xml을 찾을 수 없습니다")
            xml = zf.read("content.xml").decode("utf-8", errors="ignore")
    except zipfile.BadZipFile as exc:
        raise AttachmentError(
            f"{filename}: 유효하지 않은 OpenDocument 파일입니다"
        ) from exc

    try:
        root = SafeET.fromstring(xml)
    except SafeET.ParseError as exc:
        raise AttachmentError(f"{filename}: OpenDocument XML 파싱 실패") from exc

    parts: list[str] = []
    for elem in root.iter():
        local = _local_name(elem.tag)
        if local in {"h", "p", "span", "list-item", "table-cell"}:
            text = " ".join(t.strip() for t in elem.itertext() if t and t.strip())
            if text:
                parts.append(text)

    result = "\n".join(parts)
    if not result:
        raise AttachmentError(f"{filename}: 텍스트를 추출할 수 없습니다")
    return result[:MAX_CHARS_PER_FILE]


def _local_name(tag: str) -> str:
    return tag.split("}", 1)[-1] if "}" in tag else tag
