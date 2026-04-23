"""app/services/attachment_service.py — Extract plain text from uploaded files.

Supported formats:
- .txt / .md  : read as-is (UTF-8, latin-1 fallback)
- .pdf        : pdfplumber — text layer + table extraction
- .docx       : python-docx — paragraphs + tables in document order
- .hwp / .hwpx: ZIP+XML parsing (HWPX format)
- .xlsx / .xls: openpyxl — sheet/row/cell extraction
- .csv        : stdlib csv module
- .pptx       : python-pptx — slide title + body text extraction
- .json / .yaml / .xml / .html / .rtf : structured text normalization
- .odt / .ods / .odp                   : OpenDocument ZIP+XML extraction
- .zip                                 : supported documents inside ZIP

Per-file cap: MAX_CHARS_PER_FILE (12 000 chars / ~3 000 tokens).
Global cap:   MAX_TOTAL_CHARS   (20 000 chars) across all files in one call.
File size:    MAX_FILE_SIZE_BYTES (20 MB).
"""
from __future__ import annotations

import csv
import html
import io
import json
import logging
import re
import zipfile
from pathlib import Path
from typing import Any

from defusedxml import ElementTree as SafeET

_log = logging.getLogger("decisiondoc.attachment")

MAX_CHARS_PER_FILE  = 12_000   # ~3 000 tokens per file
MAX_CHARS           = MAX_CHARS_PER_FILE   # backward-compat alias
MAX_TOTAL_CHARS     = 20_000   # hard cap across all files in one request
MAX_FILE_SIZE_BYTES = 20 * 1024 * 1024  # 20 MB

ALLOWED_EXTENSIONS = {
    ".txt", ".md",
    ".pdf",
    ".docx",
    ".pptx",
    ".hwp", ".hwpx",
    ".xlsx", ".xls",
    ".csv", ".tsv",
    ".json", ".jsonl", ".ndjson",
    ".yaml", ".yml",
    ".xml",
    ".html", ".htm",
    ".rtf",
    ".log", ".ini", ".cfg", ".conf",
    ".odt", ".ods", ".odp",
    ".zip",
}

PLAIN_TEXT_EXTENSIONS = {
    ".txt", ".md", ".log", ".ini", ".cfg", ".conf", ".jsonl", ".ndjson",
}
IMAGE_EXTENSIONS = {
    ".png", ".jpg", ".jpeg", ".webp", ".gif", ".bmp", ".tif", ".tiff", ".heic", ".svg",
}
LEGACY_CONVERSION_HINTS = {
    ".doc": "구형 Word(.doc)는 직접 추출하지 못합니다. DOCX 또는 PDF로 변환해 업로드하세요.",
    ".ppt": "구형 PowerPoint(.ppt)는 직접 추출하지 못합니다. PPTX 또는 PDF로 변환해 업로드하세요.",
    ".xlsb": "바이너리 Excel(.xlsb)는 직접 추출하지 못합니다. XLSX, CSV 또는 PDF로 변환해 업로드하세요.",
    ".pages": "Apple Pages 문서는 직접 추출하지 못합니다. PDF 또는 DOCX로 변환해 업로드하세요.",
    ".numbers": "Apple Numbers 문서는 직접 추출하지 못합니다. XLSX 또는 CSV로 변환해 업로드하세요.",
    ".key": "Apple Keynote 문서는 직접 추출하지 못합니다. PPTX 또는 PDF로 변환해 업로드하세요.",
}
MAX_ZIP_MEMBERS = 20


class AttachmentError(Exception):
    """Raised when a file cannot be read or its format is unsupported."""


# ── Public API ────────────────────────────────────────────────────────────────

def extract_text(filename: str, raw: bytes) -> str:
    """Extract plain text from *raw* bytes.

    Args:
        filename: Original filename — used to detect the extension.
        raw:      Raw file bytes.

    Returns:
        Extracted text, truncated to MAX_CHARS_PER_FILE.

    Raises:
        AttachmentError: on unsupported format, missing optional library,
                         parse failure, or file-size limit exceeded.
    """
    if len(raw) > MAX_FILE_SIZE_BYTES:
        mb = len(raw) // (1024 * 1024)
        raise AttachmentError(
            f"{filename}: 파일 크기 {mb} MB 초과 (최대 20 MB)"
        )

    return _extract_by_extension(filename, raw, allow_zip=True)


def extract_text_with_ai_fallback(
    filename: str,
    raw: bytes,
    *,
    provider: Any | None = None,
    request_id: str = "",
) -> str:
    """Extract text locally first, then fall back to provider OCR/vision when needed.

    The local parser remains the source of truth for structured document formats.
    Provider fallback is only attempted for image files or PDFs that look like
    scanned documents with no readable text layer.
    """
    try:
        return extract_text(filename, raw)
    except AttachmentError as exc:
        if provider is None or not _should_try_provider_fallback(filename, exc):
            raise
        try:
            text = provider.extract_attachment_text(filename, raw, request_id=request_id)
        except Exception as provider_exc:
            raise AttachmentError(
                f"{filename}: 이미지/PDF OCR·비전 추출 실패 ({provider_exc})"
            ) from provider_exc
        if not text.strip():
            raise AttachmentError(f"{filename}: OCR/비전 추출 결과가 비어 있습니다")
        return text[:MAX_CHARS_PER_FILE]


def _should_try_provider_fallback(filename: str, exc: AttachmentError) -> bool:
    ext = Path(filename).suffix.lower()
    if ext in IMAGE_EXTENSIONS:
        return True
    if ext == ".pdf" and _is_scanned_pdf_error(exc):
        return True
    return False


def _is_scanned_pdf_error(exc: AttachmentError) -> bool:
    message = str(exc)
    return "스캔 이미지 PDF일 수 있습니다" in message


def _extract_by_extension(filename: str, raw: bytes, *, allow_zip: bool) -> str:
    ext = Path(filename).suffix.lower()
    if ext in IMAGE_EXTENSIONS:
        raise AttachmentError(
            f"{filename}: 이미지 파일은 아직 본문 OCR/비전 분석을 지원하지 않습니다. "
            "이미지를 PDF/DOCX/PPTX에 포함하거나 설명 텍스트를 함께 업로드하세요."
        )
    if ext in LEGACY_CONVERSION_HINTS:
        raise AttachmentError(f"{filename}: {LEGACY_CONVERSION_HINTS[ext]}")
    if ext not in ALLOWED_EXTENSIONS:
        raise AttachmentError(
            f"지원하지 않는 파일 형식입니다: '{ext}'. "
            f"지원 형식: {', '.join(sorted(ALLOWED_EXTENSIONS))}"
        )

    try:
        if ext == ".pdf":
            return _extract_pdf(raw, filename)
        if ext == ".docx":
            return _extract_docx(raw, filename)
        if ext == ".pptx":
            return _extract_pptx(raw, filename)
        if ext in (".hwp", ".hwpx"):
            return _extract_hwpx(raw, filename)
        if ext in (".xlsx", ".xls"):
            return _extract_excel(raw, filename)
        if ext in (".csv", ".tsv"):
            return _extract_csv(raw, delimiter="\t" if ext == ".tsv" else ",")
        if ext == ".json":
            return _extract_json(raw, filename)
        if ext in (".yaml", ".yml"):
            return _extract_yaml(raw, filename)
        if ext == ".xml":
            return _extract_xml(raw, filename)
        if ext in (".html", ".htm"):
            return _extract_html(raw)
        if ext == ".rtf":
            return _extract_rtf(raw)
        if ext in PLAIN_TEXT_EXTENSIONS:
            return _extract_plain(raw)
        if ext in (".odt", ".ods", ".odp"):
            return _extract_opendocument(raw, filename)
        if ext == ".zip":
            if not allow_zip:
                raise AttachmentError(f"{filename}: ZIP 내부의 중첩 ZIP은 처리하지 않습니다")
            return _extract_zip(raw, filename)
        return _extract_plain(raw)
    except AttachmentError:
        raise
    except Exception as exc:
        _log.error("[Attachment] Failed to parse %s: %s", filename, exc)
        raise AttachmentError(f"{filename} 파싱 실패: {exc}") from exc


def extract_multiple(
    files: list[tuple[str, bytes]],
    *,
    provider: Any | None = None,
    request_id: str = "",
) -> str:
    """Extract text from multiple files with a global character cap.

    Files that fail to parse emit a warning line instead of raising so that
    a single bad file does not block the rest.

    Args:
        files: List of ``(filename, raw_bytes)`` pairs.

    Returns:
        Structured text blocks separated by ``\\n\\n---\\n\\n``.
    """
    parts: list[str] = []
    total = 0

    for filename, raw in files:
        try:
            text = extract_text_with_ai_fallback(
                filename,
                raw,
                provider=provider,
                request_id=request_id,
            )
        except AttachmentError as exc:
            parts.append(f"[첨부파일: {filename}]\n⚠️ {exc}")
            _log.warning("[Attachment] %s", exc)
            continue

        remaining = MAX_TOTAL_CHARS - total
        if remaining <= 0:
            _log.warning(
                "[Attachment] Skipping %s: total char limit reached", filename
            )
            break
        if len(text) > remaining:
            if remaining > 500:
                text = text[:remaining] + "\n...(이하 생략)"
            else:
                _log.warning(
                    "[Attachment] Skipping %s: total char limit reached", filename
                )
                break

        parts.append(f"[첨부파일: {filename}]\n{text}")
        total += len(text)

    return "\n\n---\n\n".join(parts)


# ── Internal helpers ──────────────────────────────────────────────────────────

def _extract_pdf(raw: bytes, filename: str) -> str:
    """Extract structured text from PDF using pdfplumber.

    Uses per-character font metadata when available to mark headings and bold
    text. Falls back to ``extract_text()`` when char-level data is unavailable.
    Tables are preserved as ``[표]\\n...`` blocks. Section boundaries are marked
    with ``---`` dividers.
    """
    try:
        import pdfplumber
    except ImportError as exc:
        raise AttachmentError(
            "PDF 파싱 라이브러리가 설치되어 있지 않습니다. "
            "'pip install pdfplumber'를 실행하세요."
        ) from exc

    sections: list[str] = []
    with pdfplumber.open(io.BytesIO(raw)) as pdf:
        prev_had_content = False
        for page in pdf.pages:
            # ── Tables ────────────────────────────────────────────────────────
            for table in page.extract_tables() or []:
                if not table:
                    continue
                rows = []
                for row in table:
                    cells = [str(c or "").strip() for c in row]
                    if any(cells):
                        rows.append(" | ".join(cells))
                if rows:
                    sections.append("[표]\n" + "\n".join(rows))

            # ── Char-level structured extraction ──────────────────────────────
            try:
                chars = page.chars
                if not chars:
                    raise ValueError("no chars")

                # Compute average font size for the page
                sizes = [ch.get("size", 0) for ch in chars if ch.get("size")]
                avg_size = sum(sizes) / len(sizes) if sizes else 0.0

                # Group chars into lines by their top coordinate (y)
                lines: dict[float, list[dict]] = {}
                for ch in chars:
                    y = round(ch.get("top", 0), 1)
                    lines.setdefault(y, []).append(ch)

                page_parts: list[str] = []
                for y_key in sorted(lines):
                    line_chars = sorted(lines[y_key], key=lambda c: c.get("x0", 0))
                    line_text = _reconstruct_pdf_line_text(line_chars)
                    if not line_text:
                        continue

                    # Determine dominant font size + bold for the line
                    line_sizes = [c.get("size", 0) for c in line_chars if c.get("size")]
                    line_avg_size = sum(line_sizes) / len(line_sizes) if line_sizes else 0.0
                    is_bold = any("Bold" in (c.get("fontname") or "") for c in line_chars)
                    is_heading = avg_size > 0 and line_avg_size > avg_size * 1.2

                    if is_heading:
                        page_parts.append(f"## {line_text}")
                    elif is_bold:
                        page_parts.append(f"**{line_text}**")
                    else:
                        page_parts.append(line_text)

                if page_parts:
                    page_text = "\n".join(page_parts)
                    if prev_had_content:
                        sections.append("---")
                    sections.append(page_text)
                    prev_had_content = True

            except Exception:
                # Fallback: use plain extract_text()
                text = page.extract_text(x_tolerance=3, y_tolerance=3)
                if text and text.strip():
                    if prev_had_content:
                        sections.append("---")
                    sections.append(text.strip())
                    prev_had_content = True

    result = "\n\n".join(sections)
    if not result.strip():
        raise AttachmentError(
            f"{filename}: 텍스트를 추출할 수 없습니다 "
            "(스캔 이미지 PDF일 수 있습니다)"
        )
    return result[:MAX_CHARS_PER_FILE]


def extract_pdf_structured(raw: bytes, filename: str) -> dict:
    """Extract structured content from a PDF file.

    Uses pdfplumber char-level data to detect headings and build a section
    hierarchy. Falls back gracefully when char-level extraction is unavailable.

    Args:
        raw:      Raw PDF bytes.
        filename: Original filename (used in error messages).

    Returns:
        A dict with the following keys:

        * ``title``      — First detected heading (str, empty if none found).
        * ``sections``   — List of ``{"heading": str, "content": str}`` dicts.
        * ``raw_text``   — Full extracted text (plain, no markdown markup).
        * ``page_count`` — Number of pages in the PDF.
        * ``has_tables`` — True if at least one table was found.
        * ``pages``      — List of per-page dicts with text/headings/table flags.

    Raises:
        AttachmentError: When pdfplumber is not installed or the file cannot
                         be parsed.
    """
    try:
        import pdfplumber
    except ImportError as exc:
        raise AttachmentError(
            "PDF 파싱 라이브러리가 설치되어 있지 않습니다. "
            "'pip install pdfplumber'를 실행하세요."
        ) from exc

    page_count = 0
    has_tables = False
    raw_text_parts: list[str] = []
    page_summaries: list[dict[str, Any]] = []

    # Sections built incrementally: current heading + accumulated body lines
    detected_sections: list[dict] = []
    current_heading: str = ""
    current_body_lines: list[str] = []

    def _flush_section() -> None:
        nonlocal current_heading, current_body_lines
        content = "\n".join(current_body_lines).strip()
        if current_heading or content:
            detected_sections.append({"heading": current_heading, "content": content})
        current_heading = ""
        current_body_lines = []

    with pdfplumber.open(io.BytesIO(raw)) as pdf:
        page_count = len(pdf.pages)
        for page_index, page in enumerate(pdf.pages, start=1):
            page_line_texts: list[str] = []
            page_headings: list[str] = []
            # Check for tables
            tables = page.extract_tables() or []
            page_has_tables = bool(tables)
            if page_has_tables:
                has_tables = True
            for table in tables:
                if not table:
                    continue
                rows = []
                for row in table:
                    cells = [str(c or "").strip() for c in row]
                    if any(cells):
                        rows.append(" | ".join(cells))
                if rows:
                    raw_text_parts.append("\n".join(rows))
                    page_line_texts.extend(rows)

            # Char-level extraction for heading detection
            try:
                chars = page.chars
                if not chars:
                    raise ValueError("no chars")

                sizes = [ch.get("size", 0) for ch in chars if ch.get("size")]
                avg_size = sum(sizes) / len(sizes) if sizes else 0.0

                lines: dict[float, list[dict]] = {}
                for ch in chars:
                    y = round(ch.get("top", 0), 1)
                    lines.setdefault(y, []).append(ch)

                for y_key in sorted(lines):
                    line_chars = sorted(lines[y_key], key=lambda c: c.get("x0", 0))
                    line_text = _reconstruct_pdf_line_text(line_chars)
                    if not line_text:
                        continue

                    raw_text_parts.append(line_text)
                    page_line_texts.append(line_text)

                    line_sizes = [c.get("size", 0) for c in line_chars if c.get("size")]
                    line_avg_size = sum(line_sizes) / len(line_sizes) if line_sizes else 0.0
                    is_heading = avg_size > 0 and line_avg_size > avg_size * 1.2

                    if is_heading:
                        _flush_section()
                        current_heading = line_text
                        page_headings.append(line_text)
                    else:
                        current_body_lines.append(line_text)

            except Exception:
                # Fallback to plain text
                text = page.extract_text(x_tolerance=3, y_tolerance=3)
                if text and text.strip():
                    raw_text_parts.append(text.strip())
                    current_body_lines.append(text.strip())
                    page_line_texts.extend(line.strip() for line in text.splitlines() if line.strip())

            page_preview = " ".join(page_line_texts[:8]).strip()
            page_summaries.append(
                {
                    "page": page_index,
                    "headings": page_headings[:6],
                    "preview": page_preview[:300],
                    "has_tables": page_has_tables,
                }
            )

    _flush_section()

    raw_text = "\n".join(raw_text_parts).strip()

    # Derive title: first heading, or first non-empty line of raw_text
    title = ""
    if detected_sections and detected_sections[0]["heading"]:
        title = detected_sections[0]["heading"]
    elif raw_text:
        title = raw_text.splitlines()[0].strip()

    return {
        "title": title,
        "sections": detected_sections,
        "raw_text": raw_text,
        "page_count": page_count,
        "has_tables": has_tables,
        "pages": page_summaries,
    }


def _extract_docx(raw: bytes, filename: str) -> str:
    """Extract text from DOCX including tables, preserving document order."""
    try:
        from docx import Document
        from docx.table import Table
        from docx.text.paragraph import Paragraph
    except ImportError as exc:
        raise AttachmentError(
            "DOCX 파싱 라이브러리가 설치되어 있지 않습니다. "
            "'pip install python-docx'를 실행하세요."
        ) from exc

    doc = Document(io.BytesIO(raw))
    sections: list[str] = []

    for block in doc.element.body:
        tag = block.tag.split("}")[-1] if "}" in block.tag else block.tag

        if tag == "p":
            para = Paragraph(block, doc)
            text = para.text.strip()
            if text:
                sections.append(text)

        elif tag == "tbl":
            table = Table(block, doc)
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                if any(cells):
                    rows.append(" | ".join(cells))
            if rows:
                sections.append("[표]\n" + "\n".join(rows))

    return "\n".join(sections)[:MAX_CHARS_PER_FILE]


def _extract_hwpx(raw: bytes, filename: str) -> str:
    """Extract text from HWPX (ZIP + XML format)."""
    try:
        with zipfile.ZipFile(io.BytesIO(raw)) as zf:
            section_files = sorted(
                f for f in zf.namelist()
                if re.match(r"Contents/section\d+\.xml", f)
            )
            if not section_files:
                raise AttachmentError(
                    f"{filename}: HWPX 섹션 파일을 찾을 수 없습니다"
                )
            texts: list[str] = []
            for sf in section_files:
                xml = zf.read(sf).decode("utf-8", errors="ignore")
                # <hh:t> or <t> tags carry the visible text
                matches = re.findall(
                    r"<(?:hh:)?t[^>]*>([^<]+)</(?:hh:)?t>", xml
                )
                texts.extend(m.strip() for m in matches if m.strip())

    except zipfile.BadZipFile as exc:
        raise AttachmentError(
            f"{filename}: 구형 바이너리 HWP는 직접 추출하지 못합니다. "
            "HWPX, PDF 또는 DOCX로 변환해 업로드하세요."
        ) from exc

    result = "\n".join(texts)
    if not result:
        raise AttachmentError(f"{filename}: 텍스트를 추출할 수 없습니다")
    return result[:MAX_CHARS_PER_FILE]


def _extract_excel(raw: bytes, filename: str) -> str:
    """Extract text from Excel (.xlsx / .xls) via openpyxl."""
    try:
        import openpyxl
    except ImportError as exc:
        raise AttachmentError(
            "Excel 파싱 라이브러리가 설치되어 있지 않습니다. "
            "'pip install openpyxl'를 실행하세요."
        ) from exc

    wb = openpyxl.load_workbook(
        io.BytesIO(raw), read_only=True, data_only=True
    )
    sheets: list[str] = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows: list[str] = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c).strip() if c is not None else "" for c in row]
            if any(c for c in cells if c):
                rows.append(" | ".join(cells))
        if rows:
            sheets.append(f"[시트: {sheet_name}]\n" + "\n".join(rows))
    wb.close()

    return "\n\n".join(sheets)[:MAX_CHARS_PER_FILE]


def _extract_csv(raw: bytes, *, delimiter: str) -> str:
    """Extract text from CSV."""
    text = raw.decode("utf-8", errors="ignore")
    reader = csv.reader(io.StringIO(text), delimiter=delimiter)
    rows = [" | ".join(row) for row in reader if any(row)]
    return "\n".join(rows)[:MAX_CHARS_PER_FILE]


def _extract_pptx(raw: bytes, filename: str) -> str:
    """Extract text from PPTX — slide title + body text in slide order."""
    try:
        from pptx import Presentation
        from pptx.util import Pt  # noqa: F401 — confirms pptx is available
    except ImportError as exc:
        raise AttachmentError(
            "PPTX 파싱 라이브러리가 설치되어 있지 않습니다. "
            "'pip install python-pptx'를 실행하세요."
        ) from exc

    prs = Presentation(io.BytesIO(raw))
    slides: list[str] = []

    for i, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = "".join(run.text for run in para.runs).strip()
                if text:
                    parts.append(text)
        if parts:
            slides.append(f"[슬라이드 {i}]\n" + "\n".join(parts))

    result = "\n\n".join(slides)
    if not result.strip():
        raise AttachmentError(f"{filename}: 텍스트를 추출할 수 없습니다")
    return result[:MAX_CHARS_PER_FILE]


def _extract_plain(raw: bytes) -> str:
    """Extract plain text (.txt / .md)."""
    try:
        return raw.decode("utf-8")[:MAX_CHARS_PER_FILE]
    except UnicodeDecodeError:
        return raw.decode("latin-1", errors="replace")[:MAX_CHARS_PER_FILE]


def _extract_json(raw: bytes, filename: str) -> str:
    text = _extract_plain(raw)
    try:
        parsed = json.loads(text)
    except json.JSONDecodeError:
        return text[:MAX_CHARS_PER_FILE]
    return json.dumps(parsed, ensure_ascii=False, indent=2)[:MAX_CHARS_PER_FILE]


def _extract_yaml(raw: bytes, filename: str) -> str:
    try:
        import yaml
    except ImportError as exc:
        raise AttachmentError(
            "YAML 파싱 라이브러리가 설치되어 있지 않습니다. 'pip install PyYAML'를 실행하세요."
        ) from exc

    text = _extract_plain(raw)
    try:
        parsed = yaml.safe_load(text)
    except yaml.YAMLError:
        return text[:MAX_CHARS_PER_FILE]
    if parsed is None:
        return text[:MAX_CHARS_PER_FILE]
    return yaml.safe_dump(parsed, allow_unicode=True, sort_keys=False)[:MAX_CHARS_PER_FILE]


def _extract_xml(raw: bytes, filename: str) -> str:
    text = _extract_plain(raw)
    try:
        root = SafeET.fromstring(text)
    except SafeET.ParseError:
        return text[:MAX_CHARS_PER_FILE]

    parts: list[str] = []
    for elem in root.iter():
        if elem.text and elem.text.strip():
            parts.append(elem.text.strip())
    result = "\n".join(parts)
    if not result:
        raise AttachmentError(f"{filename}: XML에서 텍스트를 추출할 수 없습니다")
    return result[:MAX_CHARS_PER_FILE]


def _extract_html(raw: bytes) -> str:
    text = _extract_plain(raw)
    text = re.sub(r"(?is)<(script|style)\b[^>]*>.*?</\1>", " ", text)
    text = re.sub(r"(?i)<br\s*/?>", "\n", text)
    text = re.sub(r"(?i)</(p|div|li|tr|h[1-6])>", "\n", text)
    text = re.sub(r"(?s)<[^>]+>", " ", text)
    text = html.unescape(text)
    text = re.sub(r"[ \t]+\n", "\n", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    text = re.sub(r"[ \t]{2,}", " ", text)
    return text.strip()[:MAX_CHARS_PER_FILE]


def _extract_rtf(raw: bytes) -> str:
    text = raw.decode("latin-1", errors="ignore")
    text = re.sub(
        r"\\'([0-9a-fA-F]{2})",
        lambda m: bytes.fromhex(m.group(1)).decode("cp1252", errors="ignore"),
        text,
    )
    text = re.sub(r"\\u(-?\d+)\??", lambda m: _decode_rtf_unicode(m.group(1)), text)
    text = text.replace("\\par", "\n").replace("\\line", "\n").replace("\\tab", "\t")
    text = re.sub(r"\\[a-zA-Z]+-?\d* ?", "", text)
    text = text.replace("{", "").replace("}", "")
    text = re.sub(r"\n{3,}", "\n\n", text)
    text = re.sub(r"[ \t]{2,}", " ", text)
    return text.strip()[:MAX_CHARS_PER_FILE]


def _decode_rtf_unicode(value: str) -> str:
    codepoint = int(value)
    if codepoint < 0:
        codepoint += 65536
    try:
        return chr(codepoint)
    except ValueError:
        return ""


def _extract_opendocument(raw: bytes, filename: str) -> str:
    try:
        with zipfile.ZipFile(io.BytesIO(raw)) as zf:
            if "content.xml" not in zf.namelist():
                raise AttachmentError(f"{filename}: OpenDocument content.xml을 찾을 수 없습니다")
            xml = zf.read("content.xml").decode("utf-8", errors="ignore")
    except zipfile.BadZipFile as exc:
        raise AttachmentError(
            f"{filename}: 유효하지 않은 OpenDocument 파일입니다"
        ) from exc

    try:
        root = SafeET.fromstring(xml)
    except SafeET.ParseError as exc:
        raise AttachmentError(f"{filename}: OpenDocument XML 파싱 실패") from exc

    parts: list[str] = []
    for elem in root.iter():
        local = _local_name(elem.tag)
        if local in {"h", "p", "span", "list-item", "table-cell"}:
            text = " ".join(t.strip() for t in elem.itertext() if t and t.strip())
            if text:
                parts.append(text)

    result = "\n".join(parts)
    if not result:
        raise AttachmentError(f"{filename}: 텍스트를 추출할 수 없습니다")
    return result[:MAX_CHARS_PER_FILE]


def _extract_zip(raw: bytes, filename: str) -> str:
    try:
        with zipfile.ZipFile(io.BytesIO(raw)) as zf:
            infos = [info for info in zf.infolist() if not info.is_dir()]
            if not infos:
                raise AttachmentError(f"{filename}: ZIP 안에 읽을 파일이 없습니다")

            parts: list[str] = []
            for info in infos[:MAX_ZIP_MEMBERS]:
                inner_name = Path(info.filename).name or info.filename
                inner_ext = Path(inner_name).suffix.lower()
                if inner_ext == ".zip":
                    parts.append(f"[압축 내부 파일: {info.filename}]\n⚠️ 중첩 ZIP은 건너뜁니다.")
                    continue
                if inner_ext in IMAGE_EXTENSIONS:
                    parts.append(
                        f"[압축 내부 파일: {info.filename}]\n⚠️ 이미지 파일은 OCR/비전 분석을 아직 지원하지 않아 건너뜁니다."
                    )
                    continue
                if inner_ext in LEGACY_CONVERSION_HINTS:
                    parts.append(f"[압축 내부 파일: {info.filename}]\n⚠️ {LEGACY_CONVERSION_HINTS[inner_ext]}")
                    continue
                if inner_ext not in ALLOWED_EXTENSIONS:
                    parts.append(f"[압축 내부 파일: {info.filename}]\n⚠️ 지원하지 않는 형식이라 건너뜁니다.")
                    continue
                if info.file_size > MAX_FILE_SIZE_BYTES:
                    parts.append(
                        f"[압축 내부 파일: {info.filename}]\n⚠️ 파일 크기가 20 MB를 초과해 건너뜁니다."
                    )
                    continue

                inner_raw = zf.read(info)
                try:
                    text = _extract_by_extension(inner_name, inner_raw, allow_zip=False)
                except AttachmentError as exc:
                    parts.append(f"[압축 내부 파일: {info.filename}]\n⚠️ {exc}")
                    continue
                parts.append(f"[압축 내부 파일: {info.filename}]\n{text}")
    except zipfile.BadZipFile as exc:
        raise AttachmentError(f"{filename}: 유효하지 않은 ZIP 파일입니다") from exc

    result = "\n\n---\n\n".join(parts)
    if not result.strip():
        raise AttachmentError(f"{filename}: ZIP 안에서 읽을 수 있는 텍스트를 찾지 못했습니다")
    return result[:MAX_CHARS_PER_FILE]


def _local_name(tag: str) -> str:
    return tag.split("}", 1)[-1] if "}" in tag else tag


def _reconstruct_pdf_line_text(line_chars: list[dict]) -> str:
    if not line_chars:
        return ""

    pieces: list[str] = []
    prev_char: dict | None = None
    for char in line_chars:
        text = str(char.get("text", ""))
        if not text:
            continue
        if prev_char is not None and _should_insert_pdf_space(prev_char, char):
            pieces.append(" ")
        pieces.append(text)
        prev_char = char
    return "".join(pieces).strip()


def _should_insert_pdf_space(prev_char: dict, curr_char: dict) -> bool:
    prev_text = str(prev_char.get("text", ""))
    curr_text = str(curr_char.get("text", ""))
    if not prev_text or not curr_text:
        return False
    if prev_text.isspace() or curr_text.isspace():
        return False

    prev_x1 = float(prev_char.get("x1", prev_char.get("x0", 0.0) + _pdf_char_width(prev_char)))
    curr_x0 = float(curr_char.get("x0", 0.0))
    gap = curr_x0 - prev_x1
    if gap <= 0:
        return False

    width_threshold = min(_pdf_char_width(prev_char), _pdf_char_width(curr_char)) * 0.25
    absolute_threshold = 1.8
    return gap >= max(absolute_threshold, width_threshold)


def _pdf_char_width(char: dict) -> float:
    x0 = float(char.get("x0", 0.0))
    x1 = float(char.get("x1", x0))
    width = x1 - x0
    if width > 0:
        return width

    text = str(char.get("text", ""))
    size = float(char.get("size", 0.0) or 0.0)
    if not size:
        return 0.0
    if any("\uac00" <= ch <= "\ud7a3" or "\u4e00" <= ch <= "\u9fff" for ch in text):
        return size * 0.88
    if text.isalnum():
        return size * 0.58
    return size * 0.42
