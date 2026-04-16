"""pptx_service — build in-memory PPTX files from slide structures or rendered docs."""
from __future__ import annotations

from io import BytesIO
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from app.services.export_outline import presentation_points, summarize_export_docs
from app.services.export_labels import humanize_doc_type
from app.services.markdown_utils import parse_markdown_blocks

_MAX_SLIDE_LINES = 5
_MAX_CONTENT_SLIDE_LINES = 4
_MAX_CONTENT_SLIDE_LEN = 64
_MAX_TABLE_ROWS = 6
_COLOR_BG_DARK = RGBColor(27, 33, 69)
_COLOR_BG_ACCENT = RGBColor(98, 79, 255)
_COLOR_BG_SOFT = RGBColor(244, 245, 255)
_COLOR_CARD = RGBColor(255, 255, 255)
_COLOR_CARD_SOFT = RGBColor(238, 235, 255)
_COLOR_TEXT_DARK = RGBColor(31, 37, 67)
_COLOR_TEXT_MUTED = RGBColor(99, 107, 139)
_COLOR_TEXT_LIGHT = RGBColor(255, 255, 255)
_COLOR_BORDER = RGBColor(216, 220, 242)


def _clean_slide_text(text: str) -> str:
    return " ".join(str(text).replace("**", "").replace("`", "").split())


def _chunk_lines(lines: list[str], size: int = _MAX_SLIDE_LINES, max_len: int = 78) -> list[list[str]]:
    cleaned: list[str] = []
    for raw in lines:
        cleaned.extend(_expand_slide_line(raw, max_len=max_len))
    if not cleaned:
        return []
    chunk_count = max(1, (len(cleaned) + size - 1) // size)
    balanced_size = max(1, (len(cleaned) + chunk_count - 1) // chunk_count)
    return [cleaned[idx: idx + balanced_size] for idx in range(0, len(cleaned), balanced_size)]


def _expand_slide_line(text: str, max_len: int = 78) -> list[str]:
    cleaned = _clean_slide_text(text)
    if not cleaned:
        return []
    points = presentation_points(cleaned, max_len=max_len, max_points=6)
    if points:
        return points
    if len(cleaned) <= max_len:
        return [cleaned]

    words = cleaned.split()
    parts = []
    current: list[str] = []
    current_len = 0
    for word in words:
        next_len = current_len + len(word) + (1 if current else 0)
        if current and next_len > max_len:
            parts.append(" ".join(current))
            current = [word]
            current_len = len(word)
        else:
            current.append(word)
            current_len = next_len
    if current:
        parts.append(" ".join(current))
    normalized = [_clean_slide_text(part) for part in parts if _clean_slide_text(part)]
    return normalized[:6]


def _table_block_lines(block: dict[str, Any]) -> list[str]:
    headers = [str(header).strip() for header in block.get("headers", [])]
    rows = block.get("rows", []) or []
    lines: list[str] = []
    for row in rows:
        if not isinstance(row, list):
            continue
        parts: list[str] = []
        for idx, cell in enumerate(row):
            value = _clean_slide_text(cell)
            if not value:
                continue
            header = headers[idx] if idx < len(headers) else f"항목 {idx + 1}"
            parts.append(f"{header}: {value}")
        if parts:
            lines.append(" / ".join(parts))
    return lines


def _render_slide(prs: Presentation, title: str, lines: list[str], notes: str = "") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = _clean_slide_text(title) or "문서"
    _style_text_frame(
        slide.shapes.title.text_frame,
        font_size_pt=26,
        bold=True,
        color=_COLOR_TEXT_DARK,
    )

    body_tf = slide.placeholders[1].text_frame
    body_tf.clear()
    cleaned_lines = [_clean_slide_text(line) for line in lines if _clean_slide_text(line)]
    if cleaned_lines:
        body_tf.paragraphs[0].text = cleaned_lines[0]
        for line in cleaned_lines[1:]:
            body_tf.add_paragraph().text = line
    else:
        body_tf.paragraphs[0].text = "내용 없음"
    _style_text_frame(body_tf, font_size_pt=17, color=_COLOR_TEXT_DARK)

    if notes.strip():
        slide.notes_slide.notes_text_frame.text = notes.strip()


def _add_text_box(
    slide: Any,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    font_size_pt: int = 14,
    bold: bool = False,
    color: RGBColor | None = None,
    align: PP_ALIGN | None = None,
) -> None:
    tx_box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tx_box.text_frame.text = _clean_slide_text(text)
    _style_text_frame(tx_box.text_frame, font_size_pt=font_size_pt, bold=bold, color=color, align=align)


def _set_text_frame_lines(
    text_frame: Any,
    lines: list[str],
    *,
    font_size_pt: int = 14,
    bold: bool = False,
    color: RGBColor | None = None,
    align: PP_ALIGN | None = None,
) -> None:
    cleaned_lines = [_clean_slide_text(line) for line in lines if _clean_slide_text(line)]
    text_frame.clear()
    if not cleaned_lines:
        cleaned_lines = [""]
    text_frame.paragraphs[0].text = cleaned_lines[0]
    for line in cleaned_lines[1:]:
        text_frame.add_paragraph().text = line
    _style_text_frame(text_frame, font_size_pt=font_size_pt, bold=bold, color=color, align=align)


def _style_text_frame(
    text_frame: Any,
    *,
    font_size_pt: int = 20,
    bold: bool = False,
    color: RGBColor | None = None,
    align: PP_ALIGN | None = None,
) -> None:
    for paragraph in text_frame.paragraphs:
        if align is not None:
            paragraph.alignment = align
        for run in paragraph.runs:
            run.font.size = Pt(font_size_pt)
            run.font.bold = bold
            if color is not None:
                run.font.color.rgb = color


def _set_slide_background(slide: Any, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_card(
    slide: Any,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    title: str,
    body: str | list[str],
    fill_color: RGBColor = _COLOR_CARD,
    title_color: RGBColor = _COLOR_TEXT_DARK,
    body_color: RGBColor = _COLOR_TEXT_MUTED,
) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = _COLOR_BORDER
    shape.line.width = Pt(1)
    _add_text_box(
        slide,
        left=left + 0.18,
        top=top + 0.12,
        width=width - 0.36,
        height=0.35,
        text=title,
        font_size_pt=14,
        bold=True,
        color=title_color,
    )
    body_box = slide.shapes.add_textbox(
        Inches(left + 0.18),
        Inches(top + 0.48),
        Inches(width - 0.36),
        Inches(height - 0.6),
    )
    if isinstance(body, list):
        _set_text_frame_lines(
            body_box.text_frame,
            body,
            font_size_pt=11,
            color=body_color,
        )
    else:
        body_box.text_frame.text = _clean_slide_text(body)
        _style_text_frame(body_box.text_frame, font_size_pt=11, color=body_color)


def _render_section_divider(
    prs: Presentation,
    title: str,
    subtitle: str = "",
    meta_lines: list[str] | None = None,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    _set_slide_background(slide, _COLOR_BG_DARK)
    slide.shapes.title.text = _clean_slide_text(title) or "섹션"
    _style_text_frame(
        slide.shapes.title.text_frame,
        font_size_pt=28,
        bold=True,
        color=_COLOR_TEXT_LIGHT,
    )
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = _clean_slide_text(subtitle) or "핵심 내용을 이어서 설명합니다."
        _style_text_frame(
            slide.placeholders[1].text_frame,
            font_size_pt=18,
            color=_COLOR_TEXT_LIGHT,
        )
    meta_lines = [_clean_slide_text(line) for line in (meta_lines or []) if _clean_slide_text(line)]
    if meta_lines:
        _add_card(
            slide,
            left=0.8,
            top=4.2,
            width=8.0,
            height=1.4,
            title="검토 메모",
            body=meta_lines,
            fill_color=_COLOR_CARD_SOFT,
            title_color=_COLOR_TEXT_DARK,
            body_color=_COLOR_TEXT_DARK,
        )


def _render_agenda_slide(prs: Presentation, title: str, items: list[str | dict[str, str]]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _set_slide_background(slide, _COLOR_BG_SOFT)
    slide.shapes.title.text = _clean_slide_text(title) or "발표 구성"
    _style_text_frame(
        slide.shapes.title.text_frame,
        font_size_pt=26,
        bold=True,
        color=_COLOR_TEXT_DARK,
    )
    normalized_items: list[dict[str, str]] = []
    for item in items:
        if isinstance(item, dict):
            item_title = _clean_slide_text(item.get("title", ""))
            item_detail = _clean_slide_text(item.get("detail", ""))
        else:
            item_title = _clean_slide_text(item)
            item_detail = ""
        if item_title:
            normalized_items.append({"title": item_title, "detail": item_detail})
    cleaned_items = normalized_items
    if not cleaned_items:
        cleaned_items = [{"title": "목차 없음", "detail": ""}]
    two_col_split = max(1, (len(cleaned_items) + 1) // 2)
    for idx, item in enumerate(cleaned_items[:6], start=1):
        col = 0 if idx <= two_col_split else 1
        row = idx - 1 if col == 0 else idx - 1 - two_col_split
        _add_card(
            slide,
            left=0.7 + (4.5 * col),
            top=1.4 + (0.92 * row),
            width=4.0,
            height=0.9,
            title=f"{idx:02d}",
            body=[item["title"], item["detail"]] if item["detail"] else item["title"],
            fill_color=_COLOR_CARD,
            title_color=_COLOR_BG_ACCENT,
            body_color=_COLOR_TEXT_DARK,
        )


def _render_summary_slide(prs: Presentation, summaries: list[dict[str, str]]) -> None:
    if not summaries:
        return

    page_size = 4
    pages = [summaries[idx: idx + page_size] for idx in range(0, len(summaries), page_size)]
    for page_index, page in enumerate(pages, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # title only
        _set_slide_background(slide, _COLOR_BG_SOFT)
        title = "핵심 검토 포인트"
        if len(pages) > 1:
            title = f"{title} ({page_index}/{len(pages)})"
        slide.shapes.title.text = title
        _style_text_frame(
            slide.shapes.title.text_frame,
            font_size_pt=26,
            bold=True,
            color=_COLOR_TEXT_DARK,
        )

        if len(page) <= 2:
            positions = [(0.65, 1.25), (0.65, 2.95)]
            card_width = 8.5
            card_height = 1.45
        else:
            positions = [(0.65, 1.25), (5.0, 1.25), (0.65, 3.15), (5.0, 3.15)]
            card_width = 3.95
            card_height = 1.6

        for summary, (left, top) in zip(page, positions, strict=False):
            metric_items = summary.get("metric_items") or [summary["metrics"]]
            body_lines = [
                summary.get("ppt_lead") or summary["lead"],
                f"핵심 섹션: {summary['sections']}",
                f"구성 지표: {' · '.join(metric_items)}",
            ]
            _add_card(
                slide,
                left=left,
                top=top,
                width=card_width,
                height=card_height,
                title=f"문서 {summary['index']} | {summary['label']}",
                body=body_lines,
                fill_color=_COLOR_CARD,
                title_color=_COLOR_TEXT_DARK,
                body_color=_COLOR_TEXT_MUTED,
            )


def _structured_slide_summaries(slide_outline: list[dict[str, Any]]) -> list[dict[str, str]]:
    summaries: list[dict[str, str]] = []
    for idx, item in enumerate(slide_outline[:4], start=1):
        title = _clean_slide_text(item.get("title", "")) or f"슬라이드 {idx}"
        lead = _clean_slide_text(item.get("key_content", "")) or "핵심 메시지 없음"
        lines = [line.strip() for line in str(item.get("key_content", "")).splitlines() if line.strip()]
        metrics = f"핵심 항목 {max(1, len(lines))}개"
        summaries.append(
            {
                "index": f"{idx:02d}",
                "label": title,
                "lead": lead,
                "sections": title,
                "metrics": metrics,
            }
        )
    return summaries


def _render_table_slide(prs: Presentation, title: str, headers: list[str], rows: list[list[str]], subtitle: str = "") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # title only
    _set_slide_background(slide, _COLOR_BG_SOFT)
    slide.shapes.title.text = _clean_slide_text(title) or "표"
    _style_text_frame(
        slide.shapes.title.text_frame,
        font_size_pt=24,
        bold=True,
        color=_COLOR_TEXT_DARK,
    )

    if subtitle:
        tx_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.15), Inches(8.7), Inches(0.4))
        tx_box.text_frame.text = _clean_slide_text(subtitle)
        _style_text_frame(tx_box.text_frame, font_size_pt=12, bold=False, color=_COLOR_TEXT_MUTED)

    row_count = max(2, len(rows) + 1)
    col_count = max(1, len(headers))
    top = Inches(1.6 if subtitle else 1.3)
    table = slide.shapes.add_table(row_count, col_count, Inches(0.6), top, Inches(8.8), Inches(4.6)).table

    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = _clean_slide_text(header)
        cell.fill.solid()
        cell.fill.fore_color.rgb = _COLOR_BG_ACCENT
        _style_text_frame(cell.text_frame, font_size_pt=12, bold=True, color=_COLOR_TEXT_LIGHT, align=PP_ALIGN.CENTER)

    for row_idx, row in enumerate(rows, start=1):
        for col_idx in range(col_count):
            cell = table.cell(row_idx, col_idx)
            value = row[col_idx] if col_idx < len(row) else ""
            cell.text = _clean_slide_text(value)
            if row_idx % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(250, 250, 255)
            _style_text_frame(cell.text_frame, font_size_pt=11, bold=False, color=_COLOR_TEXT_DARK)


def build_pptx(
    slide_data: dict[str, Any],
    title: str,
    *,
    include_outline_overview: bool = False,
) -> bytes:
    """Build a PPTX from structured slide data used by presentation_kr."""
    prs = Presentation()
    slide_outline = slide_data.get("slide_outline", []) or []

    cover = prs.slides.add_slide(prs.slide_layouts[0])
    _set_slide_background(cover, _COLOR_BG_DARK)
    cover.shapes.title.text = title
    _style_text_frame(
        cover.shapes.title.text_frame,
        font_size_pt=30,
        bold=True,
        color=_COLOR_TEXT_LIGHT,
        align=PP_ALIGN.CENTER,
    )
    if len(cover.placeholders) > 1:
        cover.placeholders[1].text = str(slide_data.get("presentation_goal", "")).strip()
        _style_text_frame(
            cover.placeholders[1].text_frame,
            font_size_pt=18,
            color=_COLOR_TEXT_LIGHT,
            align=PP_ALIGN.CENTER,
        )
    if include_outline_overview and slide_outline:
        top_titles = " · ".join(
            _clean_slide_text(item.get("title", "")) for item in slide_outline[:3] if _clean_slide_text(item.get("title", ""))
        )
        key_point = _clean_slide_text(slide_outline[0].get("key_content", ""))
        if top_titles:
            _add_card(
                cover,
                left=0.8,
                top=4.2,
                width=8.0,
                height=0.78,
                title="핵심 구성",
                body=top_titles,
                fill_color=_COLOR_CARD_SOFT,
                title_color=_COLOR_BG_ACCENT,
                body_color=_COLOR_TEXT_DARK,
            )
        if key_point:
            _add_card(
                cover,
                left=0.8,
                top=5.05,
                width=8.0,
                height=0.78,
                title="핵심 포인트",
                body=_clean_slide_text(_expand_slide_line(key_point, max_len=64)[0]),
                fill_color=_COLOR_CARD,
                title_color=_COLOR_BG_ACCENT,
                body_color=_COLOR_TEXT_DARK,
            )

    if include_outline_overview and slide_outline:
        agenda_items = []
        for item in slide_outline:
            item_title = _clean_slide_text(item.get("title", ""))
            if not item_title:
                continue
            detail_points = presentation_points(
                str(item.get("key_content", "")),
                max_len=56,
                max_points=1,
            )
            detail = _clean_slide_text(detail_points[0]) if detail_points else ""
            agenda_items.append({"title": item_title, "detail": detail})
        _render_agenda_slide(prs, "발표 구성", agenda_items)
        _render_summary_slide(prs, _structured_slide_summaries(slide_outline))

    for item in slide_outline:
        if not isinstance(item, dict):
            continue
        raw = str(item.get("key_content", ""))
        lines = [line.strip() for line in raw.splitlines() if line.strip()]
        design_tip = str(item.get("design_tip", ""))
        _render_slide(prs, str(item.get("title", "")), lines, notes=design_tip)

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def build_pptx_from_docs(docs: list[dict[str, Any]], title: str) -> bytes:
    """Build a PPTX deck from rendered markdown docs for non-presentation bundles."""
    prs = Presentation()
    summaries = summarize_export_docs(docs)

    cover = prs.slides.add_slide(prs.slide_layouts[0])
    _set_slide_background(cover, _COLOR_BG_DARK)
    cover.shapes.title.text = title
    _style_text_frame(
        cover.shapes.title.text_frame,
        font_size_pt=30,
        bold=True,
        color=_COLOR_TEXT_LIGHT,
        align=PP_ALIGN.CENTER,
    )
    if len(cover.placeholders) > 1:
        top_labels = " · ".join(summary["label"] for summary in summaries[:3])
        subtitle = f"{len(docs)}개 문서를 발표 자료 형태로 재구성"
        if top_labels:
            subtitle = f"{subtitle}\n{top_labels}"
        cover.placeholders[1].text = subtitle
        _style_text_frame(
            cover.placeholders[1].text_frame,
            font_size_pt=18,
            color=_COLOR_TEXT_LIGHT,
            align=PP_ALIGN.CENTER,
        )
    if summaries:
        _add_card(
            cover,
            left=0.8,
            top=4.8,
            width=8.0,
            height=1.1,
            title="핵심 포인트",
            body=summaries[0].get("ppt_lead") or summaries[0]["lead"],
            fill_color=_COLOR_CARD_SOFT,
            title_color=_COLOR_BG_ACCENT,
            body_color=_COLOR_TEXT_DARK,
        )

    agenda_items = [
        {
            "title": summary["label"],
            "detail": summary.get("ppt_lead") or "",
        }
        for summary in summaries
    ]
    _render_agenda_slide(prs, "발표 구성", agenda_items)
    if summaries:
        _render_summary_slide(prs, summaries)

    for idx, doc in enumerate(docs):
        doc_markdown = str(doc.get("markdown", ""))
        doc_type = str(doc.get("doc_type", "document")).strip() or "document"
        blocks = parse_markdown_blocks(doc_markdown)
        summary = summaries[idx]
        fallback_title = humanize_doc_type(doc_type)
        current_title = fallback_title
        current_lines: list[str] = []
        current_doc_heading = fallback_title
        skip_current_section = False
        table_subtitle = _clean_slide_text(summary.get("ppt_lead") or summary["lead"])

        def flush_section() -> None:
            nonlocal current_lines
            chunks = _chunk_lines(
                current_lines,
                size=_MAX_CONTENT_SLIDE_LINES,
                max_len=_MAX_CONTENT_SLIDE_LEN,
            )
            if not chunks:
                return
            for idx, chunk in enumerate(chunks, start=1):
                slide_title = current_title if idx == 1 else f"{current_title} ({idx})"
                _render_slide(prs, slide_title, chunk)
            current_lines = []

        _render_section_divider(
            prs,
            current_doc_heading,
            summary.get("ppt_lead") or summary["lead"],
            meta_lines=[
                f"핵심 섹션: {summary['sections']}",
                f"구성 지표: {' · '.join(summary.get('metric_items') or [summary['metrics']])}",
            ],
        )

        for block in blocks:
            block_type = block.get("type")
            if block_type == "heading":
                heading_text = _clean_slide_text(block.get("text", ""))
                if heading_text:
                    flush_section()
                    if "PPT 구성 가이드" in heading_text:
                        skip_current_section = True
                        current_title = fallback_title
                        continue
                    skip_current_section = False
                    current_title = heading_text
                    if current_doc_heading == fallback_title:
                        current_doc_heading = heading_text
            elif block_type == "paragraph":
                if skip_current_section:
                    continue
                text = _clean_slide_text(block.get("text", ""))
                if text:
                    current_lines.append(text)
            elif block_type == "list_item":
                if skip_current_section:
                    continue
                text = _clean_slide_text(block.get("text", ""))
                if text:
                    current_lines.append(text)
            elif block_type == "table":
                if skip_current_section:
                    continue
                flush_section()
                headers = [str(header).strip() for header in block.get("headers", [])]
                rows = block.get("rows", []) or []
                if headers and rows:
                    for idx in range(0, len(rows), _MAX_TABLE_ROWS):
                        chunk = rows[idx: idx + _MAX_TABLE_ROWS]
                        table_title = current_title if idx == 0 else f"{current_title} ({idx // _MAX_TABLE_ROWS + 1})"
                        _render_table_slide(prs, table_title, headers, chunk, subtitle=table_subtitle)
                else:
                    current_lines.extend(_table_block_lines(block))
            elif block_type == "hr":
                flush_section()
            elif block_type == "blank" and current_lines:
                # Preserve paragraph boundaries without forcing empty bullets.
                continue

        flush_section()

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()
