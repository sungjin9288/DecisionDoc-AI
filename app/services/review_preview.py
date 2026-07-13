"""review_preview — build lightweight preview data and review dashboard HTML."""
from __future__ import annotations

import html
import re
import zipfile
from io import BytesIO
from typing import Any

from docx import Document
from pptx import Presentation


def _dedupe_lines(lines: list[str], limit: int) -> list[str]:
    deduped: list[str] = []
    seen: set[str] = set()
    for raw in lines:
        line = " ".join(str(raw).split()).strip()
        if not line or line in seen:
            continue
        deduped.append(line)
        seen.add(line)
        if len(deduped) >= limit:
            break
    return deduped


def collect_docx_preview_lines(raw: bytes, *, limit: int = 16) -> list[str]:
    document = Document(BytesIO(raw))
    return _dedupe_lines([paragraph.text for paragraph in document.paragraphs], limit)


def collect_pptx_preview_lines(raw: bytes, *, limit: int = 16) -> list[str]:
    presentation = Presentation(BytesIO(raw))
    lines: list[str] = []
    for idx, slide in enumerate(presentation.slides, start=1):
        title = ""
        if getattr(slide.shapes, "title", None) is not None and slide.shapes.title is not None:
            title = slide.shapes.title.text.strip()
        title = title or f"슬라이드 {idx}"
        lines.append(f"[{idx}] {title}")
        body_lines: list[str] = []
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            if getattr(slide.shapes, "title", None) is not None and shape == slide.shapes.title:
                continue
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if text:
                    body_lines.append(f"- {text}")
        lines.extend(body_lines[:2])
    return _dedupe_lines(lines, limit)


def collect_hwpx_preview_lines(raw: bytes, *, limit: int = 16) -> list[str]:
    with zipfile.ZipFile(BytesIO(raw), "r") as archive:
        xml = archive.read("Contents/section0.xml").decode("utf-8", errors="ignore")
    texts = re.findall(r"<(?:\w+:)?t>(.*?)</(?:\w+:)?t>", xml, flags=re.DOTALL)
    cleaned = [html.unescape(re.sub(r"\s+", " ", text)).strip() for text in texts]
    return _dedupe_lines(cleaned, limit)


def preview_export_bytes(export_format: str, raw: bytes, *, limit: int = 16) -> list[str]:
    if export_format == "docx":
        return collect_docx_preview_lines(raw, limit=limit)
    if export_format == "pptx":
        return collect_pptx_preview_lines(raw, limit=limit)
    if export_format == "hwp":
        return collect_hwpx_preview_lines(raw, limit=limit)
    return []


def _status_row(label: str, value: str, tone: str) -> str:
    return (
        "<div class='check-row'>"
        f"<dt>{html.escape(label)}</dt>"
        f"<dd><span class='status status-{tone}'>{html.escape(value)}</span></dd>"
        "</div>"
    )


def _request_rows(request: dict[str, Any]) -> str:
    labels = {
        "goal": "목표",
        "context": "입력 근거",
        "constraints": "제약 조건",
        "audience": "검토 대상",
    }
    rows = []
    for key, label in labels.items():
        value = str(request.get(key) or "").strip()
        if value:
            rows.append(f"<dt>{label}</dt><dd>{html.escape(value)}</dd>")
    return "".join(rows)


def _document_sections(
    markdown_files: dict[str, str],
    documents: dict[str, str],
) -> str:
    sections = []
    for index, (document_type, path) in enumerate(markdown_files.items()):
        markdown = documents.get(document_type, "")
        open_attribute = " open" if index == 0 else ""
        sections.append(
            f"<details class='document'{open_attribute}>"
            "<summary>"
            f"<strong>{html.escape(document_type)}</strong>"
            f"<span>{len(markdown.splitlines())} lines</span>"
            "</summary>"
            "<div class='document-body'>"
            f"<a class='file-link' href='{html.escape(path)}'>Markdown 원문</a>"
            f"<pre>{html.escape(markdown)}</pre>"
            "</div>"
            "</details>"
        )
    return "".join(sections)


def build_review_dashboard(
    *,
    generated_at: str,
    manifest: dict[str, Any],
    bundle_previews: dict[str, dict[str, list[str]]],
    bundle_documents: dict[str, dict[str, str]] | None = None,
    human_review_receipt_path: str | None = None,
) -> str:
    bundle_documents = bundle_documents or {}
    bundle_sections: list[str] = []
    for bundle_type, bundle in manifest.get("bundles", {}).items():
        exports = bundle.get("exports", {})
        markdown_docs = bundle.get("markdown_docs", {})
        preview_files = bundle.get("preview_files", {})
        quality = bundle.get("quality", {})
        numeric_review = quality.get("numeric_grounding_review", {})
        request = bundle.get("request", {})

        export_links = "".join(
            f'<a class="file-link" href="{html.escape(path)}">{label.upper()}</a>'
            for label, path in exports.items()
        )
        quality_rows = "".join(
            [
                _status_row(
                    "Schema validator",
                    "통과" if quality.get("validator_pass") else "실패",
                    "pass" if quality.get("validator_pass") else "fail",
                ),
                _status_row(
                    "Bundle lint",
                    "통과" if quality.get("lint_pass") else "실패",
                    "pass" if quality.get("lint_pass") else "fail",
                ),
                _status_row(
                    "수치 근거 확인",
                    "통과" if numeric_review.get("status") == "passed" else "검토 필요",
                    "pass" if numeric_review.get("status") == "passed" else "pending",
                ),
                _status_row(
                    "사실 근거 검토",
                    "완료" if quality.get("factual_grounding_verified") else "검토 필요",
                    "pass" if quality.get("factual_grounding_verified") else "pending",
                ),
                _status_row(
                    "사람의 시각 검토",
                    "완료" if quality.get("human_visual_review_completed") else "검토 필요",
                    "pass" if quality.get("human_visual_review_completed") else "pending",
                ),
            ]
        )

        preview_blocks: list[str] = []
        for fmt in ("docx", "pptx", "hwp"):
            preview_lines = bundle_previews.get(bundle_type, {}).get(fmt, [])
            if not preview_lines:
                continue
            preview_path = preview_files.get(fmt, "")
            items = "".join(f"<li>{html.escape(line)}</li>" for line in preview_lines)
            preview_blocks.append(
                "<div class='export-preview'>"
                f"<div class='preview-head'><strong>{fmt.upper()} preview</strong>"
                + (
                    f"<a href='{html.escape(preview_path)}'>txt</a>"
                    if preview_path
                    else ""
                )
                + "</div>"
                f"<ul>{items}</ul>"
                "</div>"
            )

        pdf_embed = ""
        if exports.get("pdf"):
            pdf_embed = (
                "<div class='export-preview'>"
                "<div class='preview-head'><strong>PDF preview</strong></div>"
                f"<iframe src='{html.escape(exports['pdf'])}#toolbar=0&navpanes=0'></iframe>"
                "</div>"
            )

        bundle_sections.append(
            "<section class='bundle'>"
            "<header class='bundle-header'>"
            "<div>"
            f"<p class='bundle-id'>{html.escape(bundle_type)}</p>"
            f"<h2>{html.escape(bundle.get('title', bundle_type))}</h2>"
            f"<p class='bundle-meta'>생성 문서 {bundle.get('doc_count', 0)}개</p>"
            "</div>"
            f"<div class='export-links'>{export_links}</div>"
            "</header>"
            "<div class='review-context'>"
            "<section class='context-section'>"
            "<h3>요청 근거</h3>"
            f"<dl class='request-data'>{_request_rows(request)}</dl>"
            "</section>"
            "<section class='context-section'>"
            "<h3>검증 상태</h3>"
            f"<dl class='check-list'>{quality_rows}</dl>"
            "</section>"
            "</div>"
            "<section class='documents'>"
            "<h3>생성 문서</h3>"
            f"{_document_sections(markdown_docs, bundle_documents.get(bundle_type, {}))}"
            "</section>"
            "<div class='export-grid'>"
            + "".join(preview_blocks)
            + pdf_embed
            + "</div>"
            "</section>"
        )

    summary = manifest.get("summary", {})
    overall_passed = manifest.get("status") == "passed"
    overall_status = "자동 검증 통과" if overall_passed else "추가 검토 필요"
    overall_tone = "pass" if overall_passed else "pending"
    receipt_link = (
        f'<a class="file-link" href="{html.escape(human_review_receipt_path)}">사람 검토 receipt</a>'
        if human_review_receipt_path
        else ""
    )

    return f"""<!doctype html>
<html lang="ko">
<head>
  <meta charset="utf-8">
  <title>완성 문서 검토</title>
  <meta name="viewport" content="width=device-width, initial-scale=1">
  <link rel="icon" href="data:,">
  <style>
    :root {{
      color-scheme: light;
      --bg: #f3f5f6;
      --surface: #ffffff;
      --surface-muted: #f7f9f9;
      --border: #d6dddf;
      --border-strong: #aebabc;
      --text: #182326;
      --muted: #5b696d;
      --accent: #0f766e;
      --accent-strong: #0b5e58;
      --pass-bg: #e9f6ef;
      --pass-text: #166534;
      --pending-bg: #fff4d6;
      --pending-text: #8a4b08;
      --fail-bg: #feecec;
      --fail-text: #b42318;
    }}
    * {{ box-sizing: border-box; }}
    body {{
      margin: 0;
      font-family: Pretendard, SUIT, -apple-system, BlinkMacSystemFont, "Segoe UI", "Apple SD Gothic Neo", sans-serif;
      background: var(--bg);
      color: var(--text);
      letter-spacing: 0;
    }}
    .shell {{
      width: min(1180px, calc(100vw - 40px));
      margin: 0 auto 64px;
    }}
    .page-header {{
      display: flex;
      justify-content: space-between;
      gap: 24px;
      align-items: flex-start;
      padding: 38px 0 24px;
      border-bottom: 1px solid var(--border-strong);
    }}
    .page-header h1 {{ margin: 0 0 8px; font-size: 30px; line-height: 1.25; }}
    .page-header p {{ margin: 0; color: var(--muted); line-height: 1.55; }}
    .header-actions {{ display: flex; gap: 12px; align-items: center; flex-wrap: wrap; justify-content: flex-end; }}
    .eyebrow {{
      margin-bottom: 7px !important;
      color: var(--accent-strong) !important;
      font-size: 12px;
      font-weight: 800;
      text-transform: uppercase;
    }}
    .status {{
      display: inline-flex;
      align-items: center;
      min-height: 28px;
      padding: 4px 9px;
      border-radius: 999px;
      font-size: 12px;
      font-weight: 800;
      white-space: nowrap;
    }}
    .status-pass {{ background: var(--pass-bg); color: var(--pass-text); }}
    .status-pending {{ background: var(--pending-bg); color: var(--pending-text); }}
    .status-fail {{ background: var(--fail-bg); color: var(--fail-text); }}
    .summary {{
      display: grid;
      grid-template-columns: repeat(4, minmax(0, 1fr));
      border-bottom: 1px solid var(--border-strong);
    }}
    .metric {{ padding: 18px 20px; border-right: 1px solid var(--border); }}
    .metric:first-child {{ padding-left: 0; }}
    .metric:last-child {{ border-right: 0; }}
    .metric span {{ display: block; color: var(--muted); font-size: 12px; margin-bottom: 5px; }}
    .metric strong {{ font-size: 20px; }}
    .scope-note {{
      margin: 18px 0 0;
      padding: 13px 15px;
      border-left: 3px solid var(--pending-text);
      background: #fffaf0;
      color: #654013;
      font-size: 13px;
      line-height: 1.55;
    }}
    .bundle {{ padding: 34px 0 8px; border-bottom: 1px solid var(--border-strong); }}
    .bundle-header {{ display: flex; justify-content: space-between; gap: 20px; align-items: flex-start; }}
    .bundle-id {{ margin: 0 0 5px; color: var(--accent-strong); font-size: 12px; font-weight: 800; }}
    .bundle h2 {{ margin: 0 0 6px; font-size: 23px; line-height: 1.35; }}
    .bundle-meta {{ margin: 0; color: var(--muted); font-size: 13px; }}
    .export-links {{ display: flex; flex-wrap: wrap; gap: 8px; justify-content: flex-end; }}
    .file-link {{ color: var(--accent-strong); font-weight: 700; font-size: 13px; text-decoration: none; }}
    .file-link:hover {{ text-decoration: underline; }}
    .review-context {{
      display: grid;
      grid-template-columns: repeat(2, minmax(0, 1fr));
      gap: 28px;
      margin-top: 24px;
    }}
    h3 {{ margin: 0 0 12px; font-size: 15px; }}
    .context-section {{ min-width: 0; }}
    dl {{ margin: 0; }}
    .request-data {{ display: grid; grid-template-columns: 92px 1fr; border-top: 1px solid var(--border); }}
    .request-data dt, .request-data dd {{ margin: 0; padding: 10px 0; border-bottom: 1px solid var(--border); line-height: 1.55; }}
    .request-data dt {{ color: var(--muted); font-size: 12px; font-weight: 700; }}
    .request-data dd {{ font-size: 13px; }}
    .check-list {{ border-top: 1px solid var(--border); }}
    .check-row {{ display: flex; justify-content: space-between; gap: 16px; align-items: center; border-bottom: 1px solid var(--border); padding: 8px 0; }}
    .check-row dt {{ color: var(--muted); font-size: 13px; }}
    .check-row dd {{ margin: 0; }}
    .documents {{ margin-top: 28px; }}
    .document {{ border-top: 1px solid var(--border); background: var(--surface); }}
    .document:last-child {{ border-bottom: 1px solid var(--border); }}
    .document summary {{
      display: flex;
      justify-content: space-between;
      gap: 16px;
      align-items: center;
      min-height: 48px;
      padding: 10px 12px;
      cursor: pointer;
      list-style: none;
    }}
    .document summary::-webkit-details-marker {{ display: none; }}
    .document summary::after {{ content: "+"; color: var(--accent-strong); font-size: 18px; font-weight: 700; }}
    .document[open] summary::after {{ content: "−"; }}
    .document summary strong {{ flex: 1; }}
    .document summary span {{ color: var(--muted); font-size: 12px; }}
    .document-body {{ padding: 0 12px 14px; }}
    pre {{
      max-height: 520px;
      overflow: auto;
      margin: 10px 0 0;
      padding: 16px;
      border: 1px solid var(--border);
      border-radius: 6px;
      background: var(--surface-muted);
      color: #243236;
      font-family: "SFMono-Regular", Consolas, "Liberation Mono", monospace;
      font-size: 12px;
      line-height: 1.65;
      white-space: pre-wrap;
      overflow-wrap: anywhere;
    }}
    .export-grid {{
      display: grid;
      grid-template-columns: repeat(2, minmax(0, 1fr));
      gap: 16px;
      margin-top: 18px;
    }}
    .export-preview {{
      border: 1px solid var(--border);
      border-radius: 8px;
      padding: 14px;
      background: var(--surface);
    }}
    .preview-head {{
      display: flex;
      justify-content: space-between;
      align-items: center;
      margin-bottom: 10px;
      color: var(--muted);
    }}
    .preview-head a {{ color: var(--accent-strong); text-decoration: none; font-size: 13px; }}
    ul {{ margin: 0; padding-left: 18px; }}
    li {{ margin: 0 0 8px; line-height: 1.5; }}
    iframe {{
      width: 100%;
      height: 560px;
      border: 0;
      border-radius: 6px;
      background: white;
    }}
    @media (max-width: 960px) {{
      .summary {{ grid-template-columns: repeat(2, minmax(0, 1fr)); }}
      .metric:nth-child(2) {{ border-right: 0; }}
      .metric:nth-child(-n+2) {{ border-bottom: 1px solid var(--border); }}
      .review-context, .export-grid {{ grid-template-columns: 1fr; }}
      iframe {{ height: 420px; }}
    }}
    @media (max-width: 640px) {{
      .shell {{ width: min(100vw - 24px, 1180px); }}
      .page-header, .bundle-header {{ flex-direction: column; }}
      .page-header {{ padding-top: 24px; }}
      .page-header h1 {{ font-size: 25px; }}
      .summary {{ grid-template-columns: 1fr 1fr; }}
      .metric {{ padding: 14px 10px; }}
      .metric:first-child {{ padding-left: 10px; }}
      .request-data {{ grid-template-columns: 1fr; }}
      .request-data dt {{ padding-bottom: 0; border-bottom: 0; }}
      .request-data dd {{ padding-top: 4px; }}
      .document summary {{ align-items: flex-start; }}
    }}
  </style>
</head>
<body>
  <main class="shell">
    <header class="page-header">
      <div>
        <p class="eyebrow">Local mock evidence</p>
        <h1>완성 문서 검토</h1>
        <p>생성 시각 {html.escape(generated_at)}</p>
      </div>
      <div class="header-actions">
        {receipt_link}
        <span class="status status-{overall_tone}">{overall_status}</span>
      </div>
    </header>
    <section class="summary" aria-label="검증 요약">
      <div class="metric"><span>Bundle</span><strong>{summary.get('bundle_count', 0)}</strong></div>
      <div class="metric"><span>문서</span><strong>{summary.get('document_count', 0)}</strong></div>
      <div class="metric"><span>Validator / Lint</span><strong>{summary.get('validator_pass_count', 0)} / {summary.get('lint_pass_count', 0)}</strong></div>
      <div class="metric"><span>미근거 단위 수치</span><strong>{summary.get('unsupported_numeric_claim_count', 0)}</strong></div>
    </section>
    <p class="scope-note">이 화면은 local mock fixture의 구조와 단위 수치 coverage를 보여줍니다. 사실성, 최신성, 사람의 시각 검토, live provider 품질은 별도 확인 대상입니다.</p>
    {''.join(bundle_sections)}
  </main>
</body>
</html>"""
