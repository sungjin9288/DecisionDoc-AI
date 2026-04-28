"""Tests for POST /generate/pptx endpoint."""
import base64
from io import BytesIO
from unittest.mock import patch

from fastapi.testclient import TestClient
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from app.services.pptx_service import _chunk_lines, _render_summary_slide

_PPTX_MAGIC = b"PK\x03\x04"  # ZIP/OOXML magic bytes — all .pptx files start with this


def _create_client(tmp_path, monkeypatch):
    monkeypatch.setenv("DECISIONDOC_PROVIDER", "mock")
    monkeypatch.setenv("DECISIONDOC_PROVIDER_GENERATION", "")
    monkeypatch.setenv("DECISIONDOC_PROVIDER_ATTACHMENT", "")
    monkeypatch.setenv("DECISIONDOC_PROVIDER_VISUAL", "")
    monkeypatch.setenv("DATA_DIR", str(tmp_path))
    monkeypatch.setenv("DECISIONDOC_ENV", "dev")
    monkeypatch.setenv("DECISIONDOC_MAINTENANCE", "0")
    monkeypatch.delenv("DECISIONDOC_API_KEY", raising=False)
    monkeypatch.delenv("DECISIONDOC_API_KEYS", raising=False)

    from app.main import create_app

    return TestClient(create_app())


def test_pptx_returns_binary_for_presentation_kr(tmp_path, monkeypatch):
    """Endpoint returns a valid PPTX binary for presentation_kr bundles."""
    client = _create_client(tmp_path, monkeypatch)
    res = client.post(
        "/generate/pptx",
        json={"title": "AI 발표", "goal": "핵심 전달", "bundle_type": "presentation_kr"},
    )
    assert res.status_code == 200
    assert "presentation" in res.headers["content-type"]
    assert "attachment" in res.headers.get("content-disposition", "")
    assert res.content[:4] == _PPTX_MAGIC


def test_pptx_returns_binary_for_general_bundle(tmp_path, monkeypatch):
    """Non-presentation bundles also export as PPTX via markdown-doc conversion."""
    client = _create_client(tmp_path, monkeypatch)
    res = client.post(
        "/generate/pptx",
        json={"title": "기술 결정", "goal": "아키텍처 선택", "bundle_type": "tech_decision"},
    )
    assert res.status_code == 200
    assert "presentation" in res.headers["content-type"]
    assert res.content[:4] == _PPTX_MAGIC


def test_pptx_default_bundle_exports_successfully(tmp_path, monkeypatch):
    """Omitting bundle_type falls back to tech_decision and still exports a deck."""
    client = _create_client(tmp_path, monkeypatch)
    res = client.post(
        "/generate/pptx",
        json={"title": "테스트 결정", "goal": "테스트 목표"},
    )
    assert res.status_code == 200
    assert res.content[:4] == _PPTX_MAGIC


def test_pptx_slide_count(tmp_path, monkeypatch):
    """Mock builder returns 5 slides → cover(1) + 5 content = 6 slides total."""
    from pptx import Presentation

    client = _create_client(tmp_path, monkeypatch)
    res = client.post(
        "/generate/pptx",
        json={"title": "슬라이드 카운트 테스트", "goal": "슬라이드 수 확인", "bundle_type": "presentation_kr"},
    )
    assert res.status_code == 200
    prs = Presentation(BytesIO(res.content))
    assert len(prs.slides) == 6  # 1 cover + 5 from mock outline


def test_pptx_general_bundle_has_multiple_slides(tmp_path, monkeypatch):
    """Structured proposal bundles should include overview slides before content."""
    from pptx import Presentation

    client = _create_client(tmp_path, monkeypatch)
    res = client.post(
        "/generate/pptx",
        json={"title": "제안서 PPT", "goal": "완성형 문서 변환", "bundle_type": "proposal_kr"},
    )
    assert res.status_code == 200
    prs = Presentation(BytesIO(res.content))
    assert len(prs.slides) >= 4
    titles = [slide.shapes.title.text for slide in prs.slides if getattr(slide.shapes, "title", None)]
    assert "발표 구성" in titles
    assert "핵심 검토 포인트" in titles
    assert not any("PPT 구성 가이드" in title for title in titles)


def test_pptx_general_bundle_uses_summary_slide_content(tmp_path, monkeypatch):
    """Structured proposal PPT export should include summary-driven review text."""
    from pptx import Presentation

    client = _create_client(tmp_path, monkeypatch)
    res = client.post(
        "/generate/pptx",
        json={"title": "제안서 PPT", "goal": "완성형 문서 변환", "bundle_type": "proposal_kr"},
    )
    assert res.status_code == 200
    prs = Presentation(BytesIO(res.content))
    slide_texts = []
    for slide in prs.slides:
        text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text.append(shape.text)
        slide_texts.append("\n".join(text))

    joined = "\n".join(slide_texts)
    assert "핵심 검토 포인트" in joined
    assert "사업 추진 배경" in joined
    assert "핵심 섹션:" in joined
    assert "시각자료" in joined


def test_pptx_performance_bundle_uses_structured_slide_outline(tmp_path, monkeypatch):
    """performance_plan_kr should prefer structured slide metadata over markdown fallback."""
    from pptx import Presentation

    client = _create_client(tmp_path, monkeypatch)
    res = client.post(
        "/generate/pptx",
        json={"title": "수행계획 발표", "goal": "표 중심 발표자료", "bundle_type": "performance_plan_kr"},
    )
    assert res.status_code == 200
    prs = Presentation(BytesIO(res.content))
    assert len(prs.slides) >= 10
    titles = [slide.shapes.title.text for slide in prs.slides if getattr(slide.shapes, "title", None)]
    assert "발표 구성" in titles
    assert "핵심 검토 포인트" in titles
    assert "WBS 및 마일스톤" in titles
    assert "리스크 매트릭스" in titles
    assert not any("PPT 구성 가이드" in title for title in titles)


def test_pptx_structured_slides_include_visual_and_layout_guidance(tmp_path, monkeypatch):
    from pptx import Presentation

    client = _create_client(tmp_path, monkeypatch)
    res = client.post(
        "/generate/pptx",
        json={"title": "제안서 PPT", "goal": "완성형 문서 변환", "bundle_type": "proposal_kr"},
    )
    assert res.status_code == 200
    prs = Presentation(BytesIO(res.content))
    first_content_slide = next(
        slide
        for slide in prs.slides
        if getattr(slide.shapes, "title", None) and slide.shapes.title.text == "사업 추진 배경"
    )
    content_text = "\n".join(shape.text for shape in first_content_slide.shapes if hasattr(shape, "text") and shape.text)
    assert "핵심 메시지" in content_text
    assert "권장 시각자료" in content_text
    assert "시각자료 배치" in content_text
    assert "권장 시각자료:" in content_text
    assert "배치 가이드:" in content_text


def test_pptx_structured_slides_include_decision_and_approval_contract():
    from app.services.pptx_service import build_pptx

    deck = build_pptx(
        {
            "presentation_goal": "단계형 보고서 승인",
            "slide_outline": [
                {
                    "page": 1,
                    "title": "교차로 안전 AI 적용 방안",
                    "core_message": "우회전 일시정지 감지와 보행자 보호를 우선 적용한다.",
                    "key_content": "AI 영상 분석으로 위험 이벤트를 조기 탐지하고 운영자 알림을 자동화한다.",
                    "evidence_points": ["사고 위험 구간", "운영자 대응 시간", "CCTV 연계 가능성"],
                    "visual_type": "의사결정 매트릭스",
                    "visual_brief": "기준별 도입 우선순위 평가표",
                    "layout_hint": "좌측 메시지, 우측 매트릭스, 하단 승인 기준",
                    "decision_question": "1차 구축 범위를 교차로 안전으로 승인할 것인가?",
                    "narrative_role": "대표 승인 전에 사업 범위를 확정하는 판단 장표",
                    "content_blocks": ["문제 정의", "도입 범위", "승인 요청"],
                    "data_needs": ["사고 위험 구간별 CCTV 현황", "우회전 일시정지 위반 데이터"],
                    "acceptance_criteria": ["도입 범위가 명확함", "근거 자료가 장표에 연결됨", "PM 승인 질문에 답함"],
                }
            ],
        },
        title="테스트 보고서",
        include_outline_overview=True,
    )

    prs = Presentation(BytesIO(deck))
    joined = "\n".join(
        shape.text
        for slide in prs.slides
        for shape in slide.shapes
        if hasattr(shape, "text") and shape.text
    )
    assert "의사결정 질문" in joined
    assert "승인 기준" in joined
    assert "장표 구성" in joined
    assert "스토리 역할" in joined
    assert "검증 필요" in joined
    assert "의사결정 매트릭스" in joined
    assert "1차 구축 범위를 교차로 안전으로 승인할 것인가?" in joined
    assert "사고 위험 구간별 CCTV 현황" in joined


def test_pptx_structured_bundle_embeds_generated_visual_assets(tmp_path, monkeypatch):
    client = _create_client(tmp_path, monkeypatch)
    fake_assets = [
        {
            "asset_id": "asset-1",
            "doc_type": "proposal_kr",
            "slide_title": "사업 추진 배경",
            "visual_type": "현장 사진",
            "visual_brief": "실행 현장과 운영 환경을 보여주는 참고 이미지",
            "layout_hint": "오른쪽 반영역 이미지",
            "source_kind": "provider_image",
            "source_model": "mock-image",
            "prompt": "prompt",
            "media_type": "image/png",
            "encoding": "base64",
            "content_base64": "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMCAO5tm8sAAAAASUVORK5CYII=",
        }
    ]

    with patch("app.routers.generate.generate_visual_assets_from_docs", return_value=fake_assets):
        res = client.post(
            "/generate/pptx",
            json={"title": "제안서 PPT", "goal": "완성형 문서 변환", "bundle_type": "proposal_kr"},
        )

    assert res.status_code == 200
    prs = Presentation(BytesIO(res.content))
    target_slide = next(
        slide
        for slide in prs.slides
        if getattr(slide.shapes, "title", None) and slide.shapes.title.text == "사업 추진 배경"
    )
    assert any(shape.shape_type == MSO_SHAPE_TYPE.PICTURE for shape in target_slide.shapes)


def test_pptx_export_uses_svg_asset_labels_as_editable_visual_evidence():
    from app.services.pptx_service import build_pptx

    svg = (
        '<svg xmlns="http://www.w3.org/2000/svg" width="1200" height="675">'
        '<text x="100" y="100">추진 일정</text>'
        '<text x="140" y="260">착수</text>'
        '<text x="420" y="260">중간 점검</text>'
        '<text x="700" y="260">최종 승인</text>'
        "</svg>"
    )
    deck = build_pptx(
        {
            "presentation_goal": "SVG asset fallback 검증",
            "slide_outline": [
                {
                    "page": 1,
                    "title": "추진 일정",
                    "core_message": "단계별 실행 일정을 승인한다.",
                    "key_content": "로드맵을 기준으로 PM과 대표 승인 게이트를 확정한다.",
                    "visual_type": "타임라인",
                    "visual_brief": "선택된 SVG 타임라인 자산",
                    "layout_hint": "우측 타임라인",
                    "decision_question": "이 일정으로 착수할 것인가?",
                    "acceptance_criteria": ["착수/점검/승인 게이트가 명확함"],
                }
            ],
        },
        title="SVG 시각자료 테스트",
        visual_assets=[
            {
                "asset_id": "svg-asset-1",
                "slide_title": "추진 일정",
                "visual_type": "타임라인",
                "visual_brief": "SVG로 생성된 로드맵",
                "layout_hint": "우측 패널",
                "media_type": "image/svg+xml",
                "encoding": "base64",
                "content_base64": base64.b64encode(svg.encode("utf-8")).decode("ascii"),
            }
        ],
    )

    prs = Presentation(BytesIO(deck))
    joined = "\n".join(
        shape.text
        for slide in prs.slides
        for shape in slide.shapes
        if hasattr(shape, "text") and shape.text
    )
    assert "타임라인 도식" in joined
    assert "착수" in joined
    assert "중간 점검" in joined
    assert "최종 승인" in joined


def test_pptx_performance_bundle_renders_timeline_and_governance_visuals(tmp_path, monkeypatch):
    from pptx import Presentation

    client = _create_client(tmp_path, monkeypatch)
    res = client.post(
        "/generate/pptx",
        json={"title": "수행계획 발표", "goal": "표 중심 발표자료", "bundle_type": "performance_plan_kr"},
    )
    assert res.status_code == 200
    prs = Presentation(BytesIO(res.content))

    timeline_slide = next(
        slide
        for slide in prs.slides
        if getattr(slide.shapes, "title", None) and slide.shapes.title.text == "WBS 및 마일스톤"
    )
    governance_slide = next(
        slide
        for slide in prs.slides
        if getattr(slide.shapes, "title", None) and slide.shapes.title.text == "추진 거버넌스"
    )

    timeline_text = "\n".join(shape.text for shape in timeline_slide.shapes if hasattr(shape, "text") and shape.text)
    governance_text = "\n".join(shape.text for shape in governance_slide.shapes if hasattr(shape, "text") and shape.text)

    assert "타임라인 도식" in timeline_text
    assert "거버넌스 구조" in governance_text


def test_pptx_summary_slide_prefers_short_presentation_lead(tmp_path, monkeypatch):
    """Summary slides should use shorter PPT-oriented lead text instead of full prose blocks."""
    from pptx import Presentation

    client = _create_client(tmp_path, monkeypatch)
    res = client.post(
        "/generate/pptx",
        json={"title": "제안서 PPT", "goal": "완성형 문서 변환", "bundle_type": "proposal_kr"},
    )
    assert res.status_code == 200
    prs = Presentation(BytesIO(res.content))
    summary_slide = next(slide for slide in prs.slides if getattr(slide.shapes, "title", None) and slide.shapes.title.text == "핵심 검토 포인트")
    summary_text = "\n".join(shape.text for shape in summary_slide.shapes if hasattr(shape, "text") and shape.text)
    assert "핵심 섹션:" in summary_text
    assert "구성 지표:" in summary_text
    assert " / 핵심 섹션:" not in summary_text
    assert "하나의 사업 범위로 묶은 안입니다. 사업 배경과 현황 문제를 정책 목표와 연결해 설명하고" not in summary_text


def test_chunk_lines_rebalances_small_tail():
    chunks = _chunk_lines(
        [
            "하나",
            "둘",
            "셋",
            "넷",
            "다섯",
            "여섯",
        ]
    )
    assert [len(chunk) for chunk in chunks] == [3, 3]


def test_chunk_lines_uses_tighter_max_len_for_dense_content():
    chunks = _chunk_lines(
        [
            "수행계획서는 계약 범위, 일정, 산출물, 투입 인력, 승인 게이트를 하나의 실행 문서로 정리한 결과물입니다."
        ],
        size=4,
        max_len=32,
    )
    flattened = [line for chunk in chunks for line in chunk]
    assert len(flattened) >= 2
    assert all(len(line) <= 32 for line in flattened)


def test_pptx_agenda_slide_includes_short_lead_detail(tmp_path, monkeypatch):
    """Agenda cards should show a short presentation-oriented detail line, not titles only."""
    from pptx import Presentation

    client = _create_client(tmp_path, monkeypatch)
    res = client.post(
        "/generate/pptx",
        json={"title": "제안서 PPT", "goal": "완성형 문서 변환", "bundle_type": "proposal_kr"},
    )
    assert res.status_code == 200
    prs = Presentation(BytesIO(res.content))
    agenda_slide = next(
        slide
        for slide in prs.slides
        if getattr(slide.shapes, "title", None) and slide.shapes.title.text == "발표 구성"
    )
    agenda_text = "\n".join(shape.text for shape in agenda_slide.shapes if hasattr(shape, "text") and shape.text)
    assert "표지" in agenda_text
    assert "사업명: 제안서 PPT" in agenda_text
    assert "정책 환경 변화" in agenda_text


def test_render_summary_slide_paginates_when_many_documents():
    prs = Presentation()
    summaries = [
        {
            "index": f"{idx:02d}",
            "label": f"문서 {idx}",
            "lead": f"문서 {idx} 요약입니다.",
            "ppt_lead": f"문서 {idx} 발표 요약",
            "sections": "요약 · 본문",
            "metrics": "표 2개",
            "metric_items": ["표 2개", "목록 1개"],
        }
        for idx in range(1, 6)
    ]

    _render_summary_slide(prs, summaries)

    titles = [slide.shapes.title.text for slide in prs.slides if getattr(slide.shapes, "title", None)]
    assert titles == ["핵심 검토 포인트 (1/2)", "핵심 검토 포인트 (2/2)"]
