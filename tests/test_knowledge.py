"""Tests for Knowledge Store and Knowledge API endpoints."""
from __future__ import annotations

import io
import json
import threading

import pytest
from fastapi.testclient import TestClient


# ── KnowledgeStore unit tests ─────────────────────────────────────────────────

class TestKnowledgeStore:
    def test_add_and_list(self, tmp_path):
        from app.storage.knowledge_store import KnowledgeStore
        store = KnowledgeStore("proj1", data_dir=str(tmp_path), tenant_id="system")
        entry = store.add_document("report.pdf", "분기 실적 보고서 내용입니다.")
        docs = store.list_documents()
        assert len(docs) == 1
        assert docs[0]["filename"] == "report.pdf"
        assert docs[0]["doc_id"] == entry.doc_id

    def test_get_document(self, tmp_path):
        from app.storage.knowledge_store import KnowledgeStore
        store = KnowledgeStore("proj1", data_dir=str(tmp_path), tenant_id="system")
        entry = store.add_document("spec.docx", "제품 스펙 문서", tags=["spec"])
        fetched = store.get_document(entry.doc_id)
        assert fetched is not None
        assert fetched.text == "제품 스펙 문서"
        assert fetched.tags == ["spec"]

    def test_delete_document(self, tmp_path):
        from app.storage.knowledge_store import KnowledgeStore
        store = KnowledgeStore("proj1", data_dir=str(tmp_path), tenant_id="system")
        entry = store.add_document("old.txt", "오래된 문서")
        assert store.delete_document(entry.doc_id) is True
        assert store.get_document(entry.doc_id) is None
        assert len(store.list_documents()) == 0

    def test_delete_nonexistent(self, tmp_path):
        from app.storage.knowledge_store import KnowledgeStore
        store = KnowledgeStore("proj1", data_dir=str(tmp_path), tenant_id="system")
        assert store.delete_document("nonexistent") is False

    def test_build_context_empty(self, tmp_path):
        from app.storage.knowledge_store import KnowledgeStore
        store = KnowledgeStore("empty_proj", data_dir=str(tmp_path), tenant_id="system")
        assert store.build_context() == ""

    def test_build_context_with_docs(self, tmp_path):
        from app.storage.knowledge_store import KnowledgeStore
        store = KnowledgeStore("proj2", data_dir=str(tmp_path), tenant_id="system")
        store.add_document("guide.md", "개발 가이드라인 문서")
        store.add_document("arch.md", "아키텍처 설계서")
        ctx = store.build_context()
        assert "프로젝트 지식" in ctx
        assert "guide.md" in ctx
        assert "arch.md" in ctx

    def test_update_style(self, tmp_path):
        from app.storage.knowledge_store import KnowledgeStore
        store = KnowledgeStore("proj3", data_dir=str(tmp_path), tenant_id="system")
        entry = store.add_document("report.txt", "보고서")
        style = {"formality": "합쇼체", "density": "간결", "sentence_endings": ["입니다", "합니다"]}
        result = store.update_style(entry.doc_id, style)
        assert result is True
        fetched = store.get_document(entry.doc_id)
        assert fetched.style_profile["formality"] == "합쇼체"

    def test_build_style_context_no_styles(self, tmp_path):
        from app.storage.knowledge_store import KnowledgeStore
        store = KnowledgeStore("proj4", data_dir=str(tmp_path), tenant_id="system")
        store.add_document("doc.txt", "내용")
        assert store.build_style_context() == ""

    def test_build_style_context_with_styles(self, tmp_path):
        from app.storage.knowledge_store import KnowledgeStore
        store = KnowledgeStore("proj5", data_dir=str(tmp_path), tenant_id="system")
        entry = store.add_document("doc.txt", "내용")
        store.update_style(entry.doc_id, {
            "formality": "합쇼체",
            "density": "상세",
            "sentence_endings": ["입니다", "합니다", "습니다"],
            "summary": "격식체 상세 문서",
        })
        ctx = store.build_style_context()
        assert "스타일 가이드" in ctx
        assert "합쇼체" in ctx

    def test_max_docs_eviction(self, tmp_path):
        from app.storage.knowledge_store import KnowledgeStore, MAX_DOCS_PER_PROJECT
        store = KnowledgeStore("proj_evict", data_dir=str(tmp_path), tenant_id="system")
        for i in range(MAX_DOCS_PER_PROJECT + 2):
            store.add_document(f"doc{i}.txt", f"내용 {i}")
        docs = store.list_documents()
        assert len(docs) == MAX_DOCS_PER_PROJECT

    def test_context_respects_max_chars(self, tmp_path):
        from app.storage.knowledge_store import KnowledgeStore
        store = KnowledgeStore("proj_big", data_dir=str(tmp_path), tenant_id="system")
        # 대용량 텍스트
        big_text = "가" * 5000
        store.add_document("big1.txt", big_text)
        store.add_document("big2.txt", big_text)
        ctx = store.build_context(max_chars=3000)
        assert len(ctx) <= 3200  # 약간의 헤더 여유

    def test_add_document_persists_learning_metadata(self, tmp_path):
        from app.storage.knowledge_store import KnowledgeStore

        store = KnowledgeStore("proj_meta", data_dir=str(tmp_path), tenant_id="system")
        entry = store.add_document(
            "award-proposal.docx",
            "수주 완료 제안서 본문",
            tags=["공공", "교통"],
            learning_mode="approved_output",
            quality_tier="gold",
            applicable_bundles=["proposal_kr", "performance_plan_kr"],
            source_organization="국토교통부",
            reference_year=2025,
            success_state="awarded",
            notes="PPT 요약과 본문 표 구조가 우수함",
        )

        fetched = store.get_document(entry.doc_id)
        assert fetched is not None
        assert fetched.learning_mode == "approved_output"
        assert fetched.quality_tier == "gold"
        assert fetched.applicable_bundles == ["proposal_kr", "performance_plan_kr"]
        assert fetched.source_organization == "국토교통부"
        assert fetched.reference_year == 2025
        assert fetched.success_state == "awarded"
        assert fetched.notes == "PPT 요약과 본문 표 구조가 우수함"
        assert fetched.knowledge_scope["project_id"] == "proj_meta"
        assert fetched.knowledge_scope["organization"] == "국토교통부"
        assert fetched.knowledge_scope["bundle_types"] == ["proposal_kr", "performance_plan_kr"]
        assert fetched.knowledge_scope["topic_tags"] == ["공공", "교통"]

    def test_update_metadata_and_rank_documents_for_context(self, tmp_path):
        from app.storage.knowledge_store import KnowledgeStore

        store = KnowledgeStore("proj_rank", data_dir=str(tmp_path), tenant_id="system")
        generic = store.add_document(
            "generic-reference.pdf",
            "일반 참고문서",
            learning_mode="reference",
            quality_tier="working",
        )
        targeted = store.add_document(
            "mobility-proposal.docx",
            "모빌리티 제안서 승인본",
            tags=["모빌리티", "제안"],
            learning_mode="approved_output",
            quality_tier="gold",
            applicable_bundles=["proposal_kr"],
            source_organization="파주시",
            reference_year=2025,
            success_state="approved",
            source_bundle_id="report_workflow",
            source_request_id="report_workflow:rw-paju-001:slides:2",
        )

        updated = store.update_metadata(
            generic.doc_id,
            quality_tier="silver",
            notes="기본 구조 참고용",
        )
        assert updated is True

        ranked = store.rank_documents_for_context(
            bundle_type="proposal_kr",
            title="파주시 모빌리티 제안",
            goal="승인 가능한 제안서 작성",
            source_organization="파주시",
            report_workflow_id="rw-paju-001",
        )
        assert ranked[0]["doc_id"] == targeted.doc_id
        assert ranked[0]["bundle_match"] is True
        assert ranked[0]["organization_match"] is True
        assert ranked[0]["report_workflow_match"] is True
        assert ranked[0]["workflow_source"] is True
        assert ranked[0]["recency_score"] > 0
        assert ranked[0]["learning_mode"] == "approved_output"
        assert ranked[0]["search_backend"] == "local_keyword"
        assert ranked[0]["query_overlap"] == len(ranked[0]["matched_query_terms"])
        assert {"모빌리티", "제안", "파주시"}.issubset(set(ranked[0]["matched_query_terms"]))
        assert "모빌리티" in ranked[0]["query_terms"]
        assert ranked[0]["knowledge_scope"]["project_id"] == "proj_rank"
        assert ranked[0]["knowledge_scope"]["report_workflow_id"] == "rw-paju-001"
        assert ranked[0]["knowledge_scope"]["bundle_types"] == ["proposal_kr"]
        assert "동일 Report Workflow 산출물" in ranked[0]["selection_reason"]
        assert "기관/고객 scope 일치" in ranked[0]["selection_reason"]
        assert "bundle `proposal_kr` 일치" in ranked[0]["selection_reason"]
        assert "graph" in ranked[0]["selection_reason"]
        assert ranked[0]["graph_relationships"]["relation_count"] >= 5
        assert "produced_by_workflow" in ranked[0]["graph_relationships"]["relation_types"]
        assert "applies_to_bundle" in ranked[0]["graph_relationships"]["relation_types"]
        assert "approved_for_reuse" in ranked[0]["graph_relationships"]["relation_types"]
        assert "workflow 관계" in ranked[0]["graph_relationship_summary"]
        assert ranked[0]["graph_relationship_score"] == 72
        assert any(item["label"] == "bundle 일치" for item in ranked[0]["score_breakdown"])
        assert any(item["label"] == "동일 workflow" for item in ranked[0]["score_breakdown"])
        assert any(item["label"] == "기관 scope 일치" for item in ranked[0]["score_breakdown"])
        assert any(item["label"] == "관계 그래프" and item["score"] == 72 for item in ranked[0]["score_breakdown"])
        assert ranked[1]["doc_id"] == generic.doc_id

        ctx = store.build_context(
            bundle_type="proposal_kr",
            title="파주시 모빌리티 제안",
            goal="승인 가능한 제안서 작성",
            source_organization="파주시",
            report_workflow_id="rw-paju-001",
        )
        assert "프로젝트 지식 학습 컨텍스트" in ctx
        assert "우선 적용 문서: proposal_kr" in ctx
        assert "선정 이유: 동일 Report Workflow 산출물" in ctx
        assert "품질 등급: gold" in ctx
        assert "출처: 파주시 / 2025" in ctx

    def test_rank_documents_can_use_sqlite_fts_backend(self, tmp_path, monkeypatch):
        from app.storage.knowledge_store import KnowledgeStore

        monkeypatch.setenv("DECISIONDOC_KNOWLEDGE_SEARCH_BACKEND", "sqlite_fts")
        store = KnowledgeStore("proj-fts", data_dir=str(tmp_path), tenant_id="system")
        entry = store.add_document(
            "smart-safety-reference.pdf",
            "본문",
            tags=["스마트", "안전"],
            source_organization="국토교통부",
            notes="스마트 안전 관제 승인본",
        )

        ranked = store.rank_documents_for_context(
            title="스마트 안전",
            source_organization="국토교통부",
        )

        assert ranked[0]["doc_id"] == entry.doc_id
        assert ranked[0]["search_backend"] == "sqlite_fts"
        assert {"스마트", "안전", "국토교통부"}.issubset(set(ranked[0]["matched_query_terms"]))

    def test_report_workflow_scope_can_prioritize_matching_approved_artifact(self, tmp_path):
        from app.storage.knowledge_store import KnowledgeStore

        store = KnowledgeStore("proj-rw-scope", data_dir=str(tmp_path), tenant_id="system")
        other_workflow = store.add_document(
            "other-approved.md",
            "다른 workflow 승인본",
            learning_mode="approved_output",
            quality_tier="gold",
            applicable_bundles=["report_workflow", "proposal_presentation"],
            source_organization="서울시",
            success_state="approved",
            source_bundle_id="report_workflow",
            source_request_id="report_workflow:rw-other:slides:1",
        )
        matching_workflow = store.add_document(
            "matching-approved.md",
            "동일 workflow 승인본",
            learning_mode="approved_output",
            quality_tier="gold",
            applicable_bundles=["report_workflow", "proposal_presentation"],
            source_organization="서울시",
            success_state="approved",
            source_bundle_id="report_workflow",
            source_request_id="report_workflow:rw-target:slides:3",
        )

        ranked = store.rank_documents_for_context(
            bundle_type="report_workflow",
            title="서울시 스마트 교통 보고서",
            goal="승인된 장표 구조 재사용",
            source_organization="서울시",
            report_workflow_id="rw-target",
        )

        assert ranked[0]["doc_id"] == matching_workflow.doc_id
        assert ranked[0]["report_workflow_match"] is True
        assert ranked[0]["workflow_source"] is True
        assert ranked[0]["organization_match"] is True
        assert ranked[0]["scope_summary"].startswith("동일 Report Workflow 산출물")
        assert ranked[1]["doc_id"] == other_workflow.doc_id
        assert ranked[1]["report_workflow_match"] is False
        assert ranked[1]["workflow_source"] is True

    def test_find_promoted_document_by_source_request_and_doc_type(self, tmp_path):
        from app.storage.knowledge_store import KnowledgeStore

        store = KnowledgeStore("proj-dedupe", data_dir=str(tmp_path), tenant_id="system")
        entry = store.add_document(
            "approved-reference.md",
            "# 승인본\n본문",
            learning_mode="approved_output",
            quality_tier="gold",
            applicable_bundles=["proposal_kr"],
            source_bundle_id="bundle-123",
            source_request_id="req-dup-001",
            source_doc_type="business_understanding",
        )

        fetched = store.find_promoted_document(
            source_request_id="req-dup-001",
            source_doc_type="business_understanding",
            source_bundle_id="bundle-123",
        )
        assert fetched is not None
        assert fetched.doc_id == entry.doc_id

    def test_same_project_id_is_isolated_by_tenant(self, tmp_path):
        from app.storage.knowledge_store import KnowledgeStore

        store_a = KnowledgeStore(
            "shared-project",
            data_dir=str(tmp_path),
            tenant_id="tenant-a",
        )
        store_b = KnowledgeStore(
            "shared-project",
            data_dir=str(tmp_path),
            tenant_id="tenant-b",
        )
        entry_a = store_a.add_document("tenant-a.txt", "tenant A private context")
        entry_b = store_b.add_document("tenant-b.txt", "tenant B private context")

        assert [item["doc_id"] for item in store_a.list_documents()] == [entry_a.doc_id]
        assert [item["doc_id"] for item in store_b.list_documents()] == [entry_b.doc_id]
        assert store_a.get_document(entry_b.doc_id) is None
        assert store_b.get_document(entry_a.doc_id) is None
        assert "tenant B private context" not in store_a.build_context()
        assert "tenant A private context" not in store_b.build_context()
        assert (
            tmp_path / "tenants" / "tenant-a" / "knowledge" / "shared-project" / "index.json"
        ).exists()

    def test_foreign_drift_is_hidden_and_preserved_during_owned_update(self, tmp_path):
        from app.storage.knowledge_store import KnowledgeStore

        store = KnowledgeStore(
            "project-a",
            data_dir=str(tmp_path),
            tenant_id="tenant-a",
        )
        owned = store.add_document("owned.txt", "owned context")
        index_path = (
            tmp_path / "tenants" / "tenant-a" / "knowledge" / "project-a" / "index.json"
        )
        records = json.loads(index_path.read_text(encoding="utf-8"))
        foreign = {
            **records[0],
            "doc_id": "abcdef123456",
            "tenant_id": "tenant-b",
            "filename": "foreign.txt",
        }
        records.append(foreign)
        index_path.write_text(json.dumps(records, ensure_ascii=False), encoding="utf-8")
        index_path.with_name("abcdef123456.txt").write_text(
            "foreign context",
            encoding="utf-8",
        )

        assert [item["doc_id"] for item in store.list_documents()] == [owned.doc_id]
        assert store.get_document("abcdef123456") is None
        assert "foreign context" not in store.build_context()
        assert store.update_metadata("abcdef123456", notes="changed") is False
        assert store.delete_document("abcdef123456") is False
        assert store.update_metadata(owned.doc_id, notes="owned update") is True

        persisted = json.loads(index_path.read_text(encoding="utf-8"))
        assert persisted[1] == foreign
        assert persisted[0]["notes"] == "owned update"

    def test_duplicate_doc_identity_fails_closed(self, tmp_path):
        from app.storage.knowledge_store import KnowledgeStore, KnowledgeStoreError

        store = KnowledgeStore(
            "project-a",
            data_dir=str(tmp_path),
            tenant_id="tenant-a",
        )
        owned = store.add_document("owned.txt", "owned context")
        index_path = (
            tmp_path / "tenants" / "tenant-a" / "knowledge" / "project-a" / "index.json"
        )
        records = json.loads(index_path.read_text(encoding="utf-8"))
        records.append({
            **records[0],
            "tenant_id": "tenant-b",
            "filename": "conflicting.txt",
        })
        index_path.write_text(json.dumps(records, ensure_ascii=False), encoding="utf-8")

        with pytest.raises(KnowledgeStoreError, match="Duplicate knowledge document identity"):
            store.list_documents()
        with pytest.raises(KnowledgeStoreError, match="Duplicate knowledge document identity"):
            store.get_document(owned.doc_id)
        with pytest.raises(KnowledgeStoreError, match="Duplicate knowledge document identity"):
            store.update_metadata(owned.doc_id, notes="changed")
        with pytest.raises(KnowledgeStoreError, match="Duplicate knowledge document identity"):
            store.delete_document(owned.doc_id)
        assert json.loads(index_path.read_text(encoding="utf-8")) == records

    def test_concurrent_instances_preserve_all_documents(self, tmp_path):
        from app.storage.knowledge_store import KnowledgeStore

        stores = [
            KnowledgeStore(
                "concurrent-project",
                data_dir=str(tmp_path),
                tenant_id="tenant-a",
            )
            for _ in range(20)
        ]
        threads = [
            threading.Thread(
                target=store.add_document,
                args=(f"document-{index}.txt", f"content {index}"),
            )
            for index, store in enumerate(stores)
        ]
        for thread in threads:
            thread.start()
        for thread in threads:
            thread.join()

        documents = stores[0].list_documents()
        assert len(documents) == 20
        assert {item["filename"] for item in documents} == {
            f"document-{index}.txt" for index in range(20)
        }

    @pytest.mark.parametrize("project_id", ["", ".", "..", "../outside", "a\\b"])
    def test_rejects_unsafe_project_storage_component(self, tmp_path, project_id):
        from app.storage.knowledge_store import KnowledgeStore

        with pytest.raises(ValueError, match="Invalid project_id"):
            KnowledgeStore(project_id, data_dir=str(tmp_path), tenant_id="tenant-a")

    def test_requires_valid_tenant_before_creating_project_paths(self, tmp_path):
        from app.storage.knowledge_store import KnowledgeStore

        with pytest.raises(TypeError):
            KnowledgeStore("project-a", data_dir=str(tmp_path))

        invalid_root = tmp_path / "invalid"
        for tenant_id in ("", " tenant-a", "tenant-a ", ".", "..", "a/b", "a\\b", "a\x00b"):
            with pytest.raises(ValueError, match="Invalid tenant_id"):
                KnowledgeStore(
                    "project-a",
                    data_dir=str(invalid_root),
                    tenant_id=tenant_id,
                )
        assert not (invalid_root / "tenants").exists()


# ── attachment_service PPTX 테스트 ────────────────────────────────────────────

class TestPptxExtraction:
    def _make_pptx(self, slides: list[list[str]]) -> bytes:
        """Helper: Create minimal PPTX bytes with given slide texts."""
        from pptx import Presentation
        from pptx.util import Inches
        buf = io.BytesIO()
        prs = Presentation()
        blank_layout = prs.slide_layouts[5]
        for texts in slides:
            slide = prs.slides.add_slide(blank_layout)
            for i, text in enumerate(texts):
                txBox = slide.shapes.add_textbox(
                    Inches(1 + i * 0.1), Inches(1), Inches(4), Inches(1)
                )
                txBox.text_frame.text = text
        prs.save(buf)
        return buf.getvalue()

    def test_extract_pptx_basic(self):
        from app.services.attachment_service import extract_text
        raw = self._make_pptx([["슬라이드 1 제목", "내용 1"], ["슬라이드 2 제목"]])
        result = extract_text("test.pptx", raw)
        assert "슬라이드 1" in result
        assert "슬라이드 1 제목" in result
        assert "슬라이드 2 제목" in result

    def test_extract_pptx_in_allowed_extensions(self):
        from app.services.attachment_service import ALLOWED_EXTENSIONS
        assert ".pptx" in ALLOWED_EXTENSIONS

    def test_extract_pptx_empty_raises(self):
        from app.services.attachment_service import extract_text, AttachmentError
        from pptx import Presentation
        buf = io.BytesIO()
        Presentation().save(buf)
        with pytest.raises(AttachmentError, match="텍스트를 추출할 수 없습니다"):
            extract_text("empty.pptx", buf.getvalue())


# ── Knowledge API endpoint 테스트 ─────────────────────────────────────────────

@pytest.fixture()
def client(tmp_path, monkeypatch):
    monkeypatch.setenv("DATA_DIR", str(tmp_path))
    monkeypatch.setenv("DECISIONDOC_PROVIDER", "mock")
    monkeypatch.setenv("DECISIONDOC_API_KEYS", "test-key")
    from app.main import create_app
    return TestClient(create_app())


HEADERS = {"X-DecisionDoc-Api-Key": "test-key"}


class TestKnowledgeAPI:
    def test_upload_txt_document(self, client, tmp_path):
        content = b"This is a test knowledge document."
        resp = client.post(
            "/knowledge/proj-api/documents",
            headers=HEADERS,
            files={"file": ("guide.txt", content, "text/plain")},
            data={
                "tags": "guide,test",
                "learning_mode": "approved_output",
                "quality_tier": "gold",
                "applicable_bundles": "proposal_kr,performance_plan_kr",
                "source_organization": "행정안전부",
                "reference_year": "2025",
                "success_state": "approved",
                "notes": "우수 제안서 레퍼런스",
            },
        )
        assert resp.status_code == 200
        body = resp.json()
        assert body["filename"] == "guide.txt"
        assert "doc_id" in body
        assert body["text_len"] > 0
        assert "guide" in body["tags"]
        assert body["learning_mode"] == "approved_output"
        assert body["quality_tier"] == "gold"
        assert body["applicable_bundles"] == ["proposal_kr", "performance_plan_kr"]
        assert body["source_organization"] == "행정안전부"
        assert body["reference_year"] == 2025
        assert body["success_state"] == "approved"
        assert body["notes"] == "우수 제안서 레퍼런스"

    def test_list_documents(self, client, tmp_path):
        content = b"Document content"
        client.post(
            "/knowledge/proj-list/documents",
            headers=HEADERS,
            files={"file": ("doc.txt", content, "text/plain")},
        )
        resp = client.get("/knowledge/proj-list/documents", headers=HEADERS)
        assert resp.status_code == 200
        body = resp.json()
        assert body["count"] == 1
        assert body["documents"][0]["filename"] == "doc.txt"

    def test_routes_keep_same_project_id_inside_request_tenant(self, client):
        client.app.state.tenant_store.create_tenant("tenant-a", "Tenant A")
        client.app.state.tenant_store.create_tenant("tenant-b", "Tenant B")

        def headers(tenant_id: str) -> dict[str, str]:
            return {
                **HEADERS,
                "X-Tenant-ID": tenant_id,
            }

        upload_a = client.post(
            "/knowledge/shared-project/documents",
            headers=headers("tenant-a"),
            files={"file": ("tenant-a.txt", b"tenant A private context", "text/plain")},
        )
        upload_b = client.post(
            "/knowledge/shared-project/documents",
            headers=headers("tenant-b"),
            files={"file": ("tenant-b.txt", b"tenant B private context", "text/plain")},
        )
        assert upload_a.status_code == upload_b.status_code == 200

        list_a = client.get(
            "/knowledge/shared-project/documents",
            headers=headers("tenant-a"),
        ).json()
        list_b = client.get(
            "/knowledge/shared-project/documents",
            headers=headers("tenant-b"),
        ).json()
        assert [item["filename"] for item in list_a["documents"]] == ["tenant-a.txt"]
        assert [item["filename"] for item in list_b["documents"]] == ["tenant-b.txt"]

        foreign_read = client.get(
            f"/knowledge/shared-project/documents/{upload_a.json()['doc_id']}",
            headers=headers("tenant-b"),
        )
        assert foreign_read.status_code == 404

    def test_get_document(self, client, tmp_path):
        content = b"Detailed content here."
        upload = client.post(
            "/knowledge/proj-get/documents",
            headers=HEADERS,
            files={"file": ("detail.txt", content, "text/plain")},
            data={"learning_mode": "template", "quality_tier": "silver"},
        )
        doc_id = upload.json()["doc_id"]
        resp = client.get(f"/knowledge/proj-get/documents/{doc_id}", headers=HEADERS)
        assert resp.status_code == 200
        assert resp.json()["text"] == "Detailed content here."
        assert resp.json()["learning_mode"] == "template"
        assert resp.json()["quality_tier"] == "silver"

    def test_get_nonexistent_document(self, client):
        resp = client.get("/knowledge/proj-x/documents/nonexistent", headers=HEADERS)
        assert resp.status_code == 404

    def test_delete_document(self, client, tmp_path):
        content = b"To be deleted"
        upload = client.post(
            "/knowledge/proj-del/documents",
            headers=HEADERS,
            files={"file": ("delete_me.txt", content, "text/plain")},
        )
        doc_id = upload.json()["doc_id"]
        resp = client.delete(f"/knowledge/proj-del/documents/{doc_id}", headers=HEADERS)
        assert resp.status_code == 200
        assert resp.json()["deleted"] is True
        # 삭제 후 조회 시 404
        resp2 = client.get(f"/knowledge/proj-del/documents/{doc_id}", headers=HEADERS)
        assert resp2.status_code == 404

    def test_context_preview(self, client, tmp_path):
        content = b"Important project knowledge"
        client.post(
            "/knowledge/proj-ctx/documents",
            headers=HEADERS,
            files={"file": ("knowledge.txt", content, "text/plain")},
            data={
                "learning_mode": "approved_output",
                "quality_tier": "gold",
                "applicable_bundles": "proposal_kr",
                "source_organization": "파주시",
                "reference_year": "2025",
                "success_state": "approved",
            },
        )
        resp = client.get(
            "/knowledge/proj-ctx/context",
            headers=HEADERS,
            params={
                "bundle_type": "proposal_kr",
                "title": "파주시 제안",
                "goal": "승인 가능한 제안서 작성",
                "source_organization": "파주시",
                "report_workflow_id": "rw-context-001",
            },
        )
        assert resp.status_code == 200
        body = resp.json()
        assert "context" in body
        assert "Important project knowledge" in body["context"]
        assert body["bundle_type"] == "proposal_kr"
        assert body["source_organization"] == "파주시"
        assert body["report_workflow_id"] == "rw-context-001"
        assert body["applied_scope"] == {
            "scope_version": "knowledge_context_preview.v1",
            "project_id": "proj-ctx",
            "bundle_type": "proposal_kr",
            "title": "파주시 제안",
            "goal": "승인 가능한 제안서 작성",
            "source_organization": "파주시",
            "report_workflow_id": "rw-context-001",
            "has_filters": True,
        }
        assert body["ranking_summary"]["total_ranked_documents"] == 1
        assert body["ranking_summary"]["returned_documents"] == 1
        assert body["ranking_summary"]["bundle_matches"] == 1
        assert body["ranking_summary"]["organization_matches"] == 1
        assert body["ranking_summary"]["report_workflow_matches"] == 0
        assert body["ranking_summary"]["graph_relationship_matches"] == 0
        assert body["ranking_summary"]["graph_relationship_score_total"] == 0
        assert body["ranking_summary"]["search_backend"] == "local_keyword"
        assert body["ranking_summary"]["top_score"] > 0
        assert body["ranking_summary"]["top_selection_reason"]
        assert body["ranking_summary"]["has_context"] is True
        assert body["ranked_documents"][0]["bundle_match"] is True
        assert body["ranked_documents"][0]["organization_match"] is True
        assert body["ranked_documents"][0]["quality_tier"] == "gold"
        assert body["ranked_documents"][0]["search_backend"] == "local_keyword"
        assert body["ranked_documents"][0]["matched_query_terms"]
        assert "selection_reason" in body["ranked_documents"][0]
        assert "graph_relationships" in body["ranked_documents"][0]
        assert "graph_relationship_summary" in body["ranked_documents"][0]
        assert "graph_relationship_score" in body["ranked_documents"][0]
        assert body["ranked_documents"][0]["score_breakdown"]
        assert body["ranked_documents"][0]["knowledge_scope"]["project_id"] == "proj-ctx"
        assert body["ranked_documents"][0]["knowledge_scope"]["organization"] == "파주시"

    def test_temporal_graph_endpoint(self, client, tmp_path):
        client.post(
            "/knowledge/proj-graph-api/documents",
            headers=HEADERS,
            files={"file": ("approved.txt", b"Approved workflow artifact", "text/plain")},
            data={
                "tags": "교통,안전",
                "learning_mode": "approved_output",
                "quality_tier": "gold",
                "applicable_bundles": "proposal_kr,report_workflow",
                "source_organization": "국토교통부",
                "success_state": "approved",
            },
        )
        resp = client.get(
            "/knowledge/proj-graph-api/temporal-graph",
            headers=HEADERS,
            params={
                "bundle_type": "proposal_kr",
                "source_organization": "국토교통부",
            },
        )
        assert resp.status_code == 200
        body = resp.json()
        assert body["graph_version"] == "knowledge_temporal_graph.v1"
        assert body["project_id"] == "proj-graph-api"
        assert body["applied_scope"] == {
            "scope_version": "knowledge_temporal_graph_scope.v1",
            "project_id": "proj-graph-api",
            "source_organization": "국토교통부",
            "report_workflow_id": "",
            "bundle_type": "proposal_kr",
            "has_filters": True,
        }
        assert body["summary"]["node_counts"]["artifact"] == 1
        assert body["summary"]["relation_counts"]["contains_artifact"] == 1
        assert body["summary"]["relation_counts"]["scoped_to_organization"] == 1
        assert body["summary"]["relation_counts"]["approved_for_reuse"] == 1

    def test_temporal_graph_export_endpoint_returns_portable_artifact(self, client, tmp_path):
        client.post(
            "/knowledge/proj-graph-export/documents",
            headers=HEADERS,
            files={"file": ("approved.txt", b"Approved workflow artifact", "text/plain")},
            data={
                "tags": "교통,안전",
                "learning_mode": "approved_output",
                "quality_tier": "gold",
                "applicable_bundles": "proposal_kr,report_workflow",
                "source_organization": "국토교통부",
                "success_state": "approved",
            },
        )

        resp = client.get(
            "/knowledge/proj-graph-export/temporal-graph/export",
            headers=HEADERS,
            params={
                "bundle_type": "proposal_kr",
                "source_organization": "국토교통부",
            },
        )

        assert resp.status_code == 200
        assert resp.headers["content-type"].startswith("application/json")
        assert resp.headers["content-disposition"] == (
            'attachment; filename="decisiondoc-knowledge-graph-proj-graph-export.json"'
        )
        body = resp.json()
        assert body["export_version"] == "decisiondoc_knowledge_graph_export.v1"
        assert body["source"] == "DecisionDoc KnowledgeStore"
        assert body["project_id"] == "proj-graph-export"
        assert body["generated_at"]
        assert body["filters"] == {
            "source_organization": "국토교통부",
            "report_workflow_id": "",
            "bundle_type": "proposal_kr",
        }
        assert body["graph"]["graph_version"] == "knowledge_temporal_graph.v1"
        assert body["graph"]["project_id"] == "proj-graph-export"
        assert body["graph"]["applied_scope"]["has_filters"] is True
        assert body["graph"]["summary"]["node_counts"]["artifact"] == 1
        assert body["graph"]["summary"]["relation_counts"]["scoped_to_organization"] == 1
        assert body["graph"]["summary"]["relation_counts"]["applies_to_bundle"] == 2

    def test_temporal_graph_export_rejects_unsupported_format(self, client):
        resp = client.get(
            "/knowledge/proj-graph-export/temporal-graph/export",
            headers=HEADERS,
            params={"format": "html"},
        )
        assert resp.status_code == 400
        assert resp.json()["detail"] == "지원하지 않는 graph export format입니다."

    def test_update_document_metadata(self, client, tmp_path):
        upload = client.post(
            "/knowledge/proj-meta/documents",
            headers=HEADERS,
            files={"file": ("meta.txt", b"metadata target", "text/plain")},
        )
        doc_id = upload.json()["doc_id"]

        resp = client.put(
            f"/knowledge/proj-meta/documents/{doc_id}/metadata",
            headers=HEADERS,
            json={
                "learning_mode": "approved_output",
                "quality_tier": "gold",
                "applicable_bundles": ["proposal_kr"],
                "source_organization": "국토교통부",
                "reference_year": 2026,
                "success_state": "awarded",
                "notes": "실수주 사례",
            },
        )
        assert resp.status_code == 200
        body = resp.json()
        assert body["updated"] is True
        assert body["learning_mode"] == "approved_output"
        assert body["quality_tier"] == "gold"
        assert body["applicable_bundles"] == ["proposal_kr"]
        assert body["source_organization"] == "국토교통부"
        assert body["reference_year"] == 2026
        assert body["success_state"] == "awarded"
        assert body["notes"] == "실수주 사례"
        assert body["knowledge_scope"]["organization"] == "국토교통부"
        assert body["knowledge_scope"]["bundle_types"] == ["proposal_kr"]

    def test_promote_generated_documents_to_knowledge(self, client, tmp_path):
        from app.storage.history_store import HistoryEntry, HistoryStore

        HistoryStore("system", base_dir=str(tmp_path)).add(
            HistoryEntry(
                entry_id="req-456",
                tenant_id="system",
                user_id="test-user",
                bundle_id="proposal_kr",
                bundle_name="proposal_kr",
                title="파주시 모빌리티 제안서",
                request_id="req-456",
                created_at="2026-04-16T00:00:00+00:00",
            )
        )

        resp = client.post(
            "/knowledge/proj-promote/promote-generated",
            headers=HEADERS,
            json={
                "title": "파주시 모빌리티 제안서",
                "bundle_type": "proposal_kr",
                "docs": [
                    {"doc_type": "business_understanding", "markdown": "# 사업 이해\n승인본 본문"},
                    {"doc_type": "execution_plan", "markdown": "# 수행 계획\n추진 전략"},
                ],
                "tags": ["공공", "교통"],
                "quality_tier": "gold",
                "success_state": "awarded",
                "source_organization": "파주시",
                "reference_year": 2026,
                "notes": "실수주 후 확정본",
                "source_bundle_id": "bundle-123",
                "source_request_id": "req-456",
            },
        )
        assert resp.status_code == 200
        body = resp.json()
        assert body["promoted"] == 2
        assert body["bundle_type"] == "proposal_kr"
        assert body["source_bundle_id"] == "bundle-123"
        assert body["source_request_id"] == "req-456"
        assert body["promoted_history_entries"] == 1
        assert body["documents"][0]["learning_mode"] == "approved_output"
        assert body["documents"][0]["quality_tier"] == "gold"
        assert body["documents"][0]["applicable_bundles"] == ["proposal_kr"]
        assert body["documents"][0]["knowledge_scope"]["project_id"] == "proj-promote"
        assert body["documents"][0]["knowledge_scope"]["bundle_types"] == ["proposal_kr"]

        history_item = HistoryStore("system", base_dir=str(tmp_path)).get_for_user("test-user")[0]
        assert history_item["knowledge_promoted"] is True
        assert history_item["knowledge_project_id"] == "proj-promote"
        assert history_item["knowledge_document_count"] == 2
        assert history_item["knowledge_quality_tier"] == "gold"
        assert history_item["knowledge_success_state"] == "awarded"
        history_detail = HistoryStore("system", base_dir=str(tmp_path)).get_entry("req-456", "test-user")
        assert history_detail is not None
        assert history_detail["knowledge_documents"][0]["doc_type"] == "business_understanding"
        assert history_detail["knowledge_documents"][1]["doc_type"] == "execution_plan"

        preview = client.get(
            "/knowledge/proj-promote/context",
            headers=HEADERS,
            params={
                "bundle_type": "proposal_kr",
                "title": "파주시 모빌리티 제안",
                "goal": "수주 가능한 제안서 작성",
            },
        )
        assert preview.status_code == 200
        ranked = preview.json()["ranked_documents"]
        assert ranked[0]["learning_mode"] == "approved_output"
        assert ranked[0]["bundle_match"] is True
        assert ranked[0]["success_state"] == "awarded"

        duplicate = client.post(
            "/knowledge/proj-promote/promote-generated",
            headers=HEADERS,
            json={
                "title": "파주시 모빌리티 제안서",
                "bundle_type": "proposal_kr",
                "docs": [
                    {"doc_type": "business_understanding", "markdown": "# 사업 이해\n승인본 본문"},
                    {"doc_type": "execution_plan", "markdown": "# 수행 계획\n추진 전략"},
                ],
                "tags": ["공공", "교통"],
                "quality_tier": "gold",
                "success_state": "awarded",
                "source_organization": "파주시",
                "reference_year": 2026,
                "notes": "실수주 후 확정본",
                "source_bundle_id": "bundle-123",
                "source_request_id": "req-456",
            },
        )
        assert duplicate.status_code == 200
        duplicate_body = duplicate.json()
        assert duplicate_body["promoted"] == 0
        assert duplicate_body["reused"] == 2
        assert duplicate_body["already_promoted"] is True

    def test_upload_requires_auth(self, client):
        resp = client.post(
            "/knowledge/proj/documents",
            files={"file": ("doc.txt", b"content", "text/plain")},
        )
        assert resp.status_code in (401, 403)

    def test_upload_unsupported_format(self, client):
        resp = client.post(
            "/knowledge/proj/documents",
            headers=HEADERS,
            files={"file": ("doc.exe", b"\x00\x01", "application/octet-stream")},
        )
        assert resp.status_code == 422
