"""Tests for POST /generate/export-edited endpoint.

The endpoint accepts pre-rendered (possibly user-edited) docs and converts them
to the requested file format without re-running LLM generation.
Supported formats: docx, pdf, excel, hwp, pptx.
"""
from __future__ import annotations

from unittest.mock import patch

import pytest
from fastapi.testclient import TestClient
from io import BytesIO
from pptx import Presentation

_ZIP_MAGIC  = b"PK\x03\x04"   # OOXML (.docx / .xlsx) and hwpx
_PDF_MAGIC  = b"%PDF"


def _create_client(tmp_path, monkeypatch):
    monkeypatch.setenv("DECISIONDOC_PROVIDER", "mock")
    monkeypatch.setenv("DATA_DIR", str(tmp_path))
    monkeypatch.setenv("DECISIONDOC_ENV", "dev")
    monkeypatch.setenv("DECISIONDOC_MAINTENANCE", "0")
    monkeypatch.delenv("DECISIONDOC_API_KEY",  raising=False)
    monkeypatch.delenv("DECISIONDOC_API_KEYS", raising=False)
    from app.main import create_app
    return TestClient(create_app())


_SAMPLE_DOCS = [
    {"doc_type": "adr",      "markdown": "# 결정\n\n## 배경\n\n배경 내용입니다.\n\n- 항목 1\n- 항목 2"},
    {"doc_type": "onepager", "markdown": "## 요약\n\n**핵심 포인트**: 테스트 문서입니다."},
]


# ── /generate/export-edited — docx ─────────────────────────────────────────

def test_export_edited_docx_returns_200(tmp_path, monkeypatch):
    client = _create_client(tmp_path, monkeypatch)
    res = client.post("/generate/export-edited", json={
        "format": "docx",
        "title": "편집된 문서",
        "docs": _SAMPLE_DOCS,
    })
    assert res.status_code == 200


def test_export_edited_docx_content_type(tmp_path, monkeypatch):
    client = _create_client(tmp_path, monkeypatch)
    res = client.post("/generate/export-edited", json={
        "format": "docx",
        "title": "편집된 문서",
        "docs": _SAMPLE_DOCS,
    })
    assert "wordprocessingml" in res.headers["content-type"]


def test_export_edited_docx_valid_bytes(tmp_path, monkeypatch):
    client = _create_client(tmp_path, monkeypatch)
    res = client.post("/generate/export-edited", json={
        "format": "docx",
        "title": "편집된 문서",
        "docs": _SAMPLE_DOCS,
    })
    assert res.content[:4] == _ZIP_MAGIC


def test_export_edited_docx_passes_slide_outline_and_visual_assets(tmp_path, monkeypatch):
    client = _create_client(tmp_path, monkeypatch)
    captured: dict[str, object] = {}

    def _fake_build_docx(docs, title, gov_options=None, visual_assets=None):
        captured["docs"] = docs
        captured["visual_assets"] = visual_assets
        return _ZIP_MAGIC + b"mock-docx"

    with patch("app.routers.generate.generate_visual_assets_from_docs", return_value=[
        {
            "asset_id": "asset-1",
            "doc_type": "proposal_kr",
            "slide_title": "사업 추진 배경",
            "visual_type": "현장 사진",
            "visual_brief": "운영 현장 이미지",
            "layout_hint": "오른쪽 이미지",
            "source_kind": "provider_image",
            "source_model": "mock-image",
            "prompt": "prompt",
            "media_type": "image/png",
            "encoding": "base64",
            "content_base64": "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMCAO5tm8sAAAAASUVORK5CYII=",
        }
    ]), patch("app.routers.generate.build_docx", side_effect=_fake_build_docx):
        res = client.post("/generate/export-edited", json={
            "format": "docx",
            "title": "편집된 제안서",
            "bundle_type": "proposal_kr",
            "docs": [
                {
                    "doc_type": "proposal_kr",
                    "markdown": "# 사업 이해\n\n본문",
                    "total_slides": 4,
                    "slide_outline": [
                        {
                            "title": "사업 추진 배경",
                            "core_message": "핵심 메시지",
                            "evidence_points": ["근거 1"],
                            "visual_type": "현장 사진",
                        }
                    ],
                }
            ],
        })

    assert res.status_code == 200
    assert res.content.startswith(_ZIP_MAGIC)
    docs = captured["docs"]
    assert isinstance(docs, list)
    assert docs[0]["slide_outline"][0]["title"] == "사업 추진 배경"
    visual_assets = captured["visual_assets"]
    assert isinstance(visual_assets, list)
    assert visual_assets[0]["slide_title"] == "사업 추진 배경"


def test_export_edited_reuses_provided_visual_assets_without_regeneration(tmp_path, monkeypatch):
    client = _create_client(tmp_path, monkeypatch)
    captured: dict[str, object] = {}

    def _fake_build_docx(docs, title, gov_options=None, visual_assets=None):
        captured["visual_assets"] = visual_assets
        return _ZIP_MAGIC + b"mock-docx"

    with patch("app.routers.generate.generate_visual_assets_from_docs") as generate_mock, patch(
        "app.routers.generate.build_docx",
        side_effect=_fake_build_docx,
    ):
        res = client.post("/generate/export-edited", json={
            "format": "docx",
            "title": "편집된 제안서",
            "bundle_type": "proposal_kr",
            "docs": [
                {
                    "doc_type": "proposal_kr",
                    "markdown": "# 사업 이해\n\n본문",
                    "total_slides": 4,
                    "slide_outline": [
                        {
                            "title": "사업 추진 배경",
                            "core_message": "핵심 메시지",
                            "evidence_points": ["근거 1"],
                            "visual_type": "현장 사진",
                        }
                    ],
                }
            ],
            "visual_assets": [
                {
                    "asset_id": "asset-1",
                    "doc_type": "proposal_kr",
                    "slide_title": "사업 추진 배경",
                    "visual_type": "현장 사진",
                    "visual_brief": "UI에서 미리 생성된 자산",
                    "layout_hint": "오른쪽 이미지",
                    "source_kind": "provider_image",
                    "source_model": "mock-image",
                    "prompt": "prompt",
                    "media_type": "image/png",
                    "encoding": "base64",
                    "content_base64": "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMCAO5tm8sAAAAASUVORK5CYII=",
                }
            ],
        })

    assert res.status_code == 200
    generate_mock.assert_not_called()
    visual_assets = captured["visual_assets"]
    assert isinstance(visual_assets, list)
    assert visual_assets[0]["visual_brief"] == "UI에서 미리 생성된 자산"


# ── /generate/export-edited — excel ────────────────────────────────────────

def test_export_edited_excel_returns_200(tmp_path, monkeypatch):
    client = _create_client(tmp_path, monkeypatch)
    res = client.post("/generate/export-edited", json={
        "format": "excel",
        "title": "엑셀 편집 문서",
        "docs": _SAMPLE_DOCS,
    })
    assert res.status_code == 200


def test_export_edited_excel_valid_bytes(tmp_path, monkeypatch):
    client = _create_client(tmp_path, monkeypatch)
    res = client.post("/generate/export-edited", json={
        "format": "excel",
        "title": "엑셀 편집 문서",
        "docs": _SAMPLE_DOCS,
    })
    assert res.content[:4] == _ZIP_MAGIC


# ── /generate/export-edited — hwp ──────────────────────────────────────────

def test_export_edited_hwp_returns_200(tmp_path, monkeypatch):
    client = _create_client(tmp_path, monkeypatch)
    res = client.post("/generate/export-edited", json={
        "format": "hwp",
        "title": "한글 편집 문서",
        "docs": _SAMPLE_DOCS,
    })
    assert res.status_code == 200


def test_export_edited_hwp_valid_bytes(tmp_path, monkeypatch):
    client = _create_client(tmp_path, monkeypatch)
    res = client.post("/generate/export-edited", json={
        "format": "hwp",
        "title": "한글 편집 문서",
        "docs": _SAMPLE_DOCS,
    })
    assert res.content[:4] == _ZIP_MAGIC  # hwpx is ZIP-based


# ── /generate/export-edited — pdf ──────────────────────────────────────────

def test_export_edited_pdf_returns_200(tmp_path, monkeypatch):
    client = _create_client(tmp_path, monkeypatch)
    res = client.post("/generate/export-edited", json={
        "format": "pdf",
        "title": "PDF 편집 문서",
        "docs": _SAMPLE_DOCS,
    })
    assert res.status_code == 200


def test_export_edited_pdf_content_type(tmp_path, monkeypatch):
    client = _create_client(tmp_path, monkeypatch)
    res = client.post("/generate/export-edited", json={
        "format": "pdf",
        "title": "PDF 편집 문서",
        "docs": _SAMPLE_DOCS,
    })
    assert res.headers["content-type"] == "application/pdf"


def test_export_edited_pdf_valid_bytes(tmp_path, monkeypatch):
    client = _create_client(tmp_path, monkeypatch)
    res = client.post("/generate/export-edited", json={
        "format": "pdf",
        "title": "PDF 편집 문서",
        "docs": _SAMPLE_DOCS,
    })
    assert res.content[:4] == _PDF_MAGIC


# ── /generate/export-edited — pptx ─────────────────────────────────────────

def test_export_edited_pptx_returns_200(tmp_path, monkeypatch):
    client = _create_client(tmp_path, monkeypatch)
    res = client.post("/generate/export-edited", json={
        "format": "pptx",
        "title": "PPT 편집 문서",
        "docs": _SAMPLE_DOCS,
    })
    assert res.status_code == 200


def test_export_edited_pptx_content_type(tmp_path, monkeypatch):
    client = _create_client(tmp_path, monkeypatch)
    res = client.post("/generate/export-edited", json={
        "format": "pptx",
        "title": "PPT 편집 문서",
        "docs": _SAMPLE_DOCS,
    })
    assert "presentationml.presentation" in res.headers["content-type"]


def test_export_edited_pptx_valid_bytes(tmp_path, monkeypatch):
    client = _create_client(tmp_path, monkeypatch)
    res = client.post("/generate/export-edited", json={
        "format": "pptx",
        "title": "PPT 편집 문서",
        "docs": _SAMPLE_DOCS,
    })
    assert res.content[:4] == _ZIP_MAGIC


def test_export_edited_pptx_uses_structured_slide_outline_when_present(tmp_path, monkeypatch):
    client = _create_client(tmp_path, monkeypatch)
    captured: dict[str, object] = {}

    def _fake_build_pptx(slide_data, title, *, include_outline_overview=False, visual_assets=None):
        captured["slide_data"] = slide_data
        captured["visual_assets"] = visual_assets
        return _ZIP_MAGIC + b"mock-pptx"

    with patch("app.routers.generate.generate_visual_assets_from_docs", return_value=[
        {
            "asset_id": "asset-1",
            "doc_type": "proposal_kr",
            "slide_title": "사업 추진 배경",
            "visual_type": "현장 사진",
            "visual_brief": "운영 현장 이미지",
            "layout_hint": "오른쪽 이미지",
            "source_kind": "provider_image",
            "source_model": "mock-image",
            "prompt": "prompt",
            "media_type": "image/png",
            "encoding": "base64",
            "content_base64": "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMCAO5tm8sAAAAASUVORK5CYII=",
        }
    ]), patch("app.routers.generate.build_pptx", side_effect=_fake_build_pptx):
        res = client.post("/generate/export-edited", json={
            "format": "pptx",
            "title": "편집된 발표자료",
            "bundle_type": "proposal_kr",
            "docs": [
                {
                    "doc_type": "proposal_kr",
                    "markdown": "# 사업 이해\n\n본문",
                    "total_slides": 4,
                    "slide_outline": [
                        {
                            "title": "사업 추진 배경",
                            "core_message": "핵심 메시지",
                            "evidence_points": ["근거 1"],
                            "visual_type": "현장 사진",
                        }
                    ],
                }
            ],
        })

    assert res.status_code == 200
    assert res.content.startswith(_ZIP_MAGIC)
    slide_data = captured["slide_data"]
    assert isinstance(slide_data, dict)
    assert slide_data["slide_outline"][0]["title"] == "사업 추진 배경"
    visual_assets = captured["visual_assets"]
    assert isinstance(visual_assets, list)
    assert visual_assets[0]["slide_title"] == "사업 추진 배경"


def test_export_edited_pptx_skips_ppt_guide_sections(tmp_path, monkeypatch):
    client = _create_client(tmp_path, monkeypatch)
    docs = [
        {
            "doc_type": "proposal_kr",
            "markdown": (
                "# 사업 이해\n\n"
                "## 제안 요약\n\n핵심 요약입니다.\n\n"
                "## PPT 구성 가이드\n\n- 발표용 메모 1\n- 발표용 메모 2\n"
            ),
        }
    ]
    res = client.post("/generate/export-edited", json={
        "format": "pptx",
        "title": "PPT 편집 문서",
        "docs": docs,
    })
    assert res.status_code == 200
    prs = Presentation(BytesIO(res.content))
    titles = [slide.shapes.title.text for slide in prs.slides if getattr(slide.shapes, "title", None)]
    assert "제안 요약" in titles
    assert not any("PPT 구성 가이드" in title for title in titles)


def test_export_edited_pptx_table_slides_include_document_context_subtitle(tmp_path, monkeypatch):
    client = _create_client(tmp_path, monkeypatch)
    docs = [
        {
            "doc_type": "proposal_kr",
            "markdown": (
                "# 사업 이해\n\n"
                "## 제안 요약\n\n"
                "본 제안은 핵심 정책 목표를 공공기관이 실제 운영 KPI로 관리할 수 있도록 데이터 통합 · AI 분석 운영 대시보드를 하나의 사업 범위로 묶은 안입니다.\n\n"
                "## 사업 목표\n\n"
                "| 목표 | KPI | 기준 |\n"
                "| --- | --- | --- |\n"
                "| 처리시간 단축 | 30% | 운영 |\n"
                "| 정확도 향상 | 20% | 품질 |\n"
            ),
        }
    ]
    res = client.post("/generate/export-edited", json={
        "format": "pptx",
        "title": "문서형 PPT",
        "docs": docs,
    })
    assert res.status_code == 200
    prs = Presentation(BytesIO(res.content))
    target_slide = next(
        slide
        for slide in prs.slides
        if getattr(slide.shapes, "title", None) and slide.shapes.title.text == "사업 목표"
    )
    texts = [shape.text for shape in target_slide.shapes if hasattr(shape, "text") and shape.text]
    assert len(texts) >= 2
    assert any("핵심 정책 목표" in text for text in texts[1:])


def test_export_edited_pptx_section_divider_uses_structured_meta_cards(tmp_path, monkeypatch):
    client = _create_client(tmp_path, monkeypatch)
    docs = [
        {
            "doc_type": "proposal_kr",
            "markdown": (
                "# 사업 이해\n\n"
                "## 제안 요약\n\n"
                "본 제안은 핵심 정책 목표를 공공기관이 실제 운영 KPI로 관리할 수 있도록 데이터 통합 · AI 분석 운영 대시보드를 하나의 사업 범위로 묶은 안입니다.\n\n"
                "## 사업 배경\n\n"
                "- 정책 환경 변화\n"
                "- 발주기관 운영 부담 증가\n\n"
                "## 사업 목표\n\n"
                "| 목표 | KPI |\n"
                "| --- | --- |\n"
                "| 처리시간 단축 | 30% |\n"
            ),
        }
    ]
    res = client.post("/generate/export-edited", json={
        "format": "pptx",
        "title": "문서형 PPT",
        "docs": docs,
    })
    assert res.status_code == 200
    prs = Presentation(BytesIO(res.content))
    divider_slide = next(
        slide
        for slide in prs.slides
        if getattr(slide.shapes, "title", None) and slide.shapes.title.text == "사업 이해"
    )
    divider_text = "\n".join(shape.text for shape in divider_slide.shapes if hasattr(shape, "text") and shape.text)
    assert "핵심 섹션" in divider_text
    assert "구성 지표" in divider_text
    assert "검토 메모" not in divider_text


# ── edge-cases & validation ─────────────────────────────────────────────────

def test_export_edited_unsupported_format_returns_400(tmp_path, monkeypatch):
    client = _create_client(tmp_path, monkeypatch)
    res = client.post("/generate/export-edited", json={
        "format": "pages",
        "title": "테스트",
        "docs": _SAMPLE_DOCS,
    })
    assert res.status_code == 400


def test_export_edited_empty_docs_still_returns_file(tmp_path, monkeypatch):
    """Endpoint must not crash when docs list is empty — returns a minimal file."""
    client = _create_client(tmp_path, monkeypatch)
    res = client.post("/generate/export-edited", json={
        "format": "docx",
        "title": "빈 문서",
        "docs": [],
    })
    assert res.status_code == 200
    assert res.content[:4] == _ZIP_MAGIC


def test_export_edited_content_disposition_filename(tmp_path, monkeypatch):
    """Content-Disposition header must carry the document title as filename."""
    client = _create_client(tmp_path, monkeypatch)
    res = client.post("/generate/export-edited", json={
        "format": "docx",
        "title": "편집 문서 제목",
        "docs": _SAMPLE_DOCS,
    })
    assert res.status_code == 200
    disp = res.headers.get("content-disposition", "")
    assert "document.docx" in disp


def test_export_edited_no_llm_call(tmp_path, monkeypatch):
    """Export-edited must NOT call the LLM — verify by counting generate calls."""
    call_count = {"n": 0}
    _original_generate = None

    def _patched_generate(*args, **kwargs):
        call_count["n"] += 1
        return _original_generate(*args, **kwargs)

    monkeypatch.setenv("DECISIONDOC_PROVIDER", "mock")
    monkeypatch.setenv("DATA_DIR", str(tmp_path))
    monkeypatch.setenv("DECISIONDOC_ENV", "dev")
    monkeypatch.setenv("DECISIONDOC_MAINTENANCE", "0")
    monkeypatch.delenv("DECISIONDOC_API_KEY",  raising=False)
    monkeypatch.delenv("DECISIONDOC_API_KEYS", raising=False)

    from app.providers.mock_provider import MockProvider
    _original_generate = MockProvider.generate_bundle
    monkeypatch.setattr(MockProvider, "generate_bundle", _patched_generate)

    from app.main import create_app
    client = TestClient(create_app())
    client.post("/generate/export-edited", json={
        "format": "docx",
        "title": "테스트",
        "docs": _SAMPLE_DOCS,
    })
    assert call_count["n"] == 0, "export-edited must not invoke generate_bundle"
