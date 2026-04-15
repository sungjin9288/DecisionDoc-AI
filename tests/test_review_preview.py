from __future__ import annotations

from io import BytesIO

from docx import Document
from pptx import Presentation
from pptx.util import Inches

from app.services.hwp_service import build_hwp
from app.services.review_preview import (
    build_review_dashboard,
    collect_docx_preview_lines,
    collect_hwpx_preview_lines,
    collect_pptx_preview_lines,
)


def _make_docx_bytes() -> bytes:
    document = Document()
    document.add_heading("제안서 요약", level=1)
    document.add_paragraph("핵심 메시지 요약 문단")
    buffer = BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def _make_pptx_bytes() -> bytes:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "발표 구성"
    slide.placeholders[1].text = "사업 개요\n추진 전략"
    buffer = BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


def test_collect_docx_preview_lines_reads_headings_and_paragraphs() -> None:
    lines = collect_docx_preview_lines(_make_docx_bytes())
    assert "제안서 요약" in lines
    assert "핵심 메시지 요약 문단" in lines


def test_collect_pptx_preview_lines_reads_slide_title_and_body() -> None:
    lines = collect_pptx_preview_lines(_make_pptx_bytes())
    assert any("발표 구성" in line for line in lines)
    assert any("사업 개요" in line for line in lines)


def test_collect_hwpx_preview_lines_reads_section_text() -> None:
    raw = build_hwp(
        docs=[
            {"doc_type": "performance_overview", "markdown": "# 수행 개요\n\n핵심 일정"},
        ],
        title="사업수행계획서",
    )
    lines = collect_hwpx_preview_lines(raw)
    assert any("사업수행계획서" in line for line in lines)
    assert any("완성형 문서 패키지" in line for line in lines)


def test_build_review_dashboard_includes_pdf_iframe_and_preview_text() -> None:
    html = build_review_dashboard(
        generated_at="20260416-101010",
        manifest={
            "generated_at": "20260416-101010",
            "bundles": {
                "proposal_kr": {
                    "title": "제안서",
                    "doc_count": 2,
                    "exports": {
                        "docx": "proposal_kr/exports/proposal_kr.docx",
                        "pdf": "proposal_kr/exports/proposal_kr.pdf",
                    },
                    "markdown_docs": {
                        "business_understanding": "proposal_kr/markdown/business_understanding.md",
                    },
                    "preview_files": {
                        "docx": "proposal_kr/previews/docx.txt",
                    },
                }
            },
        },
        bundle_previews={
            "proposal_kr": {
                "docx": ["제안서 요약", "핵심 메시지"],
            }
        },
    )
    assert "Finished Document Review" in html
    assert "iframe" in html
    assert "proposal_kr/exports/proposal_kr.pdf" in html
    assert "제안서 요약" in html
    assert "proposal_kr/previews/docx.txt" in html
